"""
肉酒場 然（ぜん）事業計画 ピッチデック生成スクリプト
python-pptx を使用 — 大フォント・余白重視・大型レイアウト版
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import io
from datetime import datetime
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

# 日本語フォント設定
_JP_FONT = "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf"
fm.fontManager.addfont(_JP_FONT)
_jp_prop = fm.FontProperties(fname=_JP_FONT)
matplotlib.rcParams["font.family"] = _jp_prop.get_name()

# === カラーパレット ===
C_BG       = RGBColor(0x1A, 0x12, 0x0B)
C_ACCENT   = RGBColor(0xC8, 0x96, 0x3C)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xF0, 0xE6, 0xD3)
C_GRAY     = RGBColor(0x8A, 0x7A, 0x6A)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN    = RGBColor(0x27, 0x7A, 0x5A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BG_HEX  = "#1A120B"
ACC_HEX = "#C8963C"
WHT_HEX = "#FFFFFF"
GRY_HEX = "#8A7A6A"
LGT_HEX = "#F0E6D3"
GRN_HEX = "#277A5A"
RED_HEX = "#C0392B"

# ─── 10年間グループ財務データ ─────────────────────────────────────

YEARS  = list(range(1, 11))

# 【肉酒場 然】年商（億円）
ZEN_REV   = [0.4, 1.7, 4.3, 13.0, 25.9, 56.2, 86.4, 103.7, 120.9, 138.2]
# 【蕎麦 然（立ち食い蕎麦）】Year4〜展開
NEW_REV   = [0.0, 0.0, 0.0,  0.5,  2.5,  6.0, 12.0,  22.0,  38.0,  55.0]
# 【SNS事業部（コンサル＋クリエイティブ）】Year3〜展開
CON_REV   = [0.0, 0.0, 0.5,  1.2,  2.0,  3.0,  4.5,   6.5,   9.0,  12.0]

# グループ合計年商
REVENUE   = [round(z+n+c, 1) for z, n, c in zip(ZEN_REV, NEW_REV, CON_REV)]

# 店舗数（然ブランド）
STORES    = [1, 2, 5, 15, 30, 65, 100, 120, 140, 160]

# 採用費（億円）：店舗拡大・本部・SNS事業部の採用コスト
RECRUIT   = [0.03, 0.05, 0.15, 0.30, 0.50, 0.80, 1.20, 1.50, 1.80, 2.00]
# 本部経費（億円）Year3〜（採用費含む）
HQ_COST   = [round(r + rec, 2) for r, rec in zip(
             [0.0, 0.0, 0.3, 0.8, 1.5, 2.5, 4.0, 5.0, 6.0, 7.0], RECRUIT)]

# 賞与引当（売上の1%相当）
BONUS     = [round(r * 0.01, 2) for r in REVENUE]

# 営業利益率（店舗レベル。本部経費・賞与控除前）
_BASE_MARGIN = [0.12, 0.13, 0.14, 0.16, 0.17, 0.18, 0.19, 0.19, 0.20, 0.20]

# 営業利益 = 売上×店舗レベル利益率 - 本部経費 - 賞与引当
PROFIT = [
    round(max(r * m - hq - bo, 0), 1)
    for r, m, hq, bo in zip(REVENUE, _BASE_MARGIN, HQ_COST, BONUS)
]
# 実効営業利益率
MARGIN_RATE = [round(p / r, 3) if r > 0 else 0 for p, r in zip(PROFIT, REVENUE)]


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color=C_ACCENT, height=Pt(1.5)):
    rect = slide.shapes.add_shape(1, left, top, width, int(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def slide_number(slide, num, total=23):
    add_textbox(slide, f"{num} / {total}",
                Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
                size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


def slide_header(sl, section_tag, title, slide_num):
    """標準ヘッダー: 左帯 + セクションタグ + タイトル + 区切り線"""
    add_rect(sl, 0, 0, Inches(0.3), SLIDE_H, C_ACCENT)
    add_textbox(sl, section_tag, Inches(0.65), Inches(0.35), Inches(6), Inches(0.45),
                size=13, bold=True, color=C_ACCENT)
    add_textbox(sl, title, Inches(0.65), Inches(0.78), Inches(12), Inches(0.75),
                size=32, bold=True, color=C_WHITE)
    add_line(sl, Inches(0.65), Inches(1.62), Inches(12.3))
    slide_number(sl, slide_num)


# ─── チャートヘルパー ─────────────────────────────────────────────

def chart_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def make_positioning_map():
    """ポジショニングマップ（散布図スタイル）"""
    fig, ax = plt.subplots(figsize=(6.8, 5.0))
    fig.patch.set_facecolor("#2A1E12")
    ax.set_facecolor("#2A1E12")

    # 競合店プロット（x=価格帯/こだわり度, y=健康・発酵度）
    competitors = [
        ("びゃく",      0.35, 0.85, GRY_HEX),
        ("なかよし",    0.55, 0.45, GRY_HEX),
        ("つむぎ堂",    0.70, 0.20, GRY_HEX),
        ("ひまり堂",    0.65, 0.15, GRY_HEX),
        ("ひまり商店",  0.75, 0.10, GRY_HEX),
    ]
    for name, x, y, c in competitors:
        ax.scatter(x, y, s=220, color=c, zorder=4, edgecolors="#1A120B", linewidths=1.5)
        ax.annotate(name, (x, y), textcoords="offset points", xytext=(8, 4),
                    color=LGT_HEX, fontsize=10.5)

    # 然（空白地帯を強調）
    ax.scatter(0.78, 0.80, s=600, color=ACC_HEX, zorder=6,
               edgecolors=WHT_HEX, linewidths=2.5, marker="*")
    ax.annotate("★ 肉酒場 然\n（空白地帯）", (0.78, 0.80),
                textcoords="offset points", xytext=(-72, 10),
                color=ACC_HEX, fontsize=12, fontweight="bold")

    # 空白地帯を丸で囲む
    from matplotlib.patches import Ellipse
    ellipse = Ellipse((0.75, 0.78), 0.38, 0.38, angle=0,
                      linewidth=1.8, edgecolor=ACC_HEX,
                      facecolor=ACC_HEX, alpha=0.08, zorder=2)
    ax.add_patch(ellipse)

    # 軸設定
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axhline(0.5, color="#4A3A28", linewidth=1.0, zorder=1)
    ax.axvline(0.5, color="#4A3A28", linewidth=1.0, zorder=1)

    # 軸ラベル
    ax.set_xlabel("大衆・低価格　　　　　　　　　上質・こだわり", color=GRY_HEX, fontsize=10.5)
    ax.set_ylabel("健康・発酵（低）                   健康・発酵（高）",
                  color=GRY_HEX, fontsize=10.5)
    ax.tick_params(left=False, bottom=False, labelleft=False, labelbottom=False)
    ax.spines[:].set_color("#3A2A18")

    fig.tight_layout(pad=0.8)
    return chart_to_image(fig)


def make_revenue_chart():
    """年商推移（3事業積み上げ棒グラフ）"""
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    fig.patch.set_facecolor(BG_HEX)
    ax.set_facecolor("#2A1E12")

    x = np.arange(len(YEARS))
    w = 0.65
    b1 = ax.bar(x, ZEN_REV, width=w, color=ACC_HEX,      label="肉酒場 然",       zorder=3)
    b2 = ax.bar(x, NEW_REV, width=w, color="#7A5C2C",     label="蕎麦 然（立ち食い）",   zorder=3,
                bottom=ZEN_REV)
    b3 = ax.bar(x, CON_REV, width=w, color=GRN_HEX,      label="コンサル事業",     zorder=3,
                bottom=[z+n for z,n in zip(ZEN_REV, NEW_REV)])

    ax.axvline(x=6, color=RED_HEX, linewidth=1.5, linestyle="--", alpha=0.85)
    ax.text(6.12, max(REVENUE) * 0.88, "100店舗\n(7年目)",
            color=RED_HEX, fontsize=9.5, va="top")

    for xi, val in zip(x, REVENUE):
        if val >= 5:
            ax.text(xi, val + 1.5, f"{val:.0f}億",
                    ha="center", va="bottom", color=WHT_HEX,
                    fontsize=9.5, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels([f"{y}年目" for y in YEARS], color=LGT_HEX, fontsize=11)
    ax.set_ylabel("グループ年商（億円）", color=GRY_HEX, fontsize=11)
    ax.tick_params(colors=GRY_HEX, labelsize=10)
    ax.spines[:].set_color("#3A2A18")
    ax.yaxis.set_tick_params(labelcolor=GRY_HEX)
    ax.grid(axis="y", color="#3A2A18", linewidth=0.6, zorder=0)
    ax.legend(loc="upper left", facecolor="#2A1E12", edgecolor="#3A2A18",
              labelcolor=LGT_HEX, fontsize=10)
    ax.set_title("グループ年商推移（10年・3事業合算）", color=WHT_HEX, fontsize=13, pad=10)
    fig.tight_layout(pad=1.2)
    return chart_to_image(fig)


def make_profit_chart():
    """営業利益推移（折れ線＋本部経費帯）"""
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    fig.patch.set_facecolor(BG_HEX)
    ax.set_facecolor("#2A1E12")

    x = np.arange(len(YEARS))

    # 本部経費エリア
    gross = [round(r * m, 1) for r, m in zip(REVENUE, _BASE_MARGIN)]
    ax.fill_between(x, gross, PROFIT, alpha=0.25, color=RED_HEX, label="本部経費・賞与引当")
    ax.fill_between(x, PROFIT, alpha=0.25, color=GRN_HEX)
    ax.plot(x, gross,  color=RED_HEX, linewidth=1.5, linestyle="--", marker="", zorder=3)
    ax.plot(x, PROFIT, color=GRN_HEX, linewidth=3.0, marker="o", markersize=8, zorder=4,
            label="営業利益（本部経費後）")

    for xi, val in zip(x, PROFIT):
        if val > 0:
            ax.text(xi, val + 0.3, f"{val:.1f}億",
                    ha="center", va="bottom", color=GRN_HEX,
                    fontsize=9.5, fontweight="bold")

    # 利益率ラベル（7・10年目）
    for yi in [6, 9]:
        rate = MARGIN_RATE[yi]
        ax.annotate(f"利益率\n{rate*100:.0f}%",
                    xy=(yi, PROFIT[yi]), xytext=(yi - 0.7, PROFIT[yi] + 1.5),
                    color=ACC_HEX, fontsize=9, arrowprops=dict(color=ACC_HEX, width=0.5, headwidth=4))

    ax.set_xticks(x)
    ax.set_xticklabels([f"{y}年目" for y in YEARS], color=LGT_HEX, fontsize=11)
    ax.set_ylabel("億円", color=GRY_HEX, fontsize=11)
    ax.tick_params(colors=GRY_HEX, labelsize=10)
    ax.spines[:].set_color("#3A2A18")
    ax.yaxis.set_tick_params(labelcolor=GRY_HEX)
    ax.grid(axis="y", color="#3A2A18", linewidth=0.6, zorder=0)
    ax.legend(loc="upper left", facecolor="#2A1E12", edgecolor="#3A2A18",
              labelcolor=LGT_HEX, fontsize=10)
    ax.set_title("営業利益推移（本部経費・賞与引当控除後）", color=WHT_HEX, fontsize=13, pad=10)
    fig.tight_layout(pad=1.2)
    return chart_to_image(fig)


# ─────────────────────────────────────────
# スライド定義
# ─────────────────────────────────────────

def s01_title(prs):
    sl = blank_slide(prs)
    bg(sl)
    add_rect(sl, 0, 0, Inches(0.3), SLIDE_H, C_ACCENT)

    add_textbox(sl, "肉酒場", Inches(1.1), Inches(0.9), Inches(7), Inches(1.0),
                size=42, bold=True, color=C_ACCENT)
    add_textbox(sl, "然（ぜん）", Inches(1.1), Inches(1.8), Inches(7), Inches(1.6),
                size=80, bold=True, color=C_WHITE)

    add_line(sl, Inches(1.1), Inches(3.6), Inches(5.5), color=C_ACCENT)
    add_textbox(sl, "糀漬け肉の定食と、和の一杯。",
                Inches(1.1), Inches(3.78), Inches(8), Inches(0.7),
                size=24, italic=True, color=C_LIGHT)

    add_textbox(sl, "事業計画書　2026年6月",
                Inches(1.1), Inches(4.7), Inches(6), Inches(0.55),
                size=16, color=C_GRAY)

    add_textbox(sl, "然", Inches(9.2), Inches(0.5), Inches(3.8), Inches(6.0),
                size=240, bold=True, color=RGBColor(0x2A, 0x1E, 0x12),
                align=PP_ALIGN.CENTER)

    slide_number(sl, 1)
    return sl


def s02_vision(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "VISION", "私たちのビジョン", 2)

    add_textbox(sl, "7年で100店舗。",
                Inches(0.65), Inches(2.0), Inches(12.5), Inches(1.3),
                size=60, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "発酵と肉の掛け合わせで、\n日本の食文化に新しい定番をつくる。",
                Inches(0.65), Inches(3.5), Inches(12.5), Inches(1.3),
                size=28, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "居酒屋×発酵×定食　という空白地帯に参入し、\n"
                "糀漬け肉料理を「日本の新定番」として全国に広める。",
                Inches(1.8), Inches(5.1), Inches(10), Inches(1.1),
                size=17, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return sl


def s03_why_now(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "MARKET", "なぜ今か　――　市場の追い風", 3)

    boxes = [
        ("外食産業", "35.7兆円", "2025年市場規模\n（前年比＋3.5%）"),
        ("居酒屋業態", "48ヶ月連続", "売上前年比プラス\nコロナ前超え達成"),
        ("発酵食品", "SNS急拡大", "健康志向×体験価値\nZ世代・40代双方に刺さる"),
    ]
    for i, (title, big, sub) in enumerate(boxes):
        x = Inches(0.65 + i * 4.2)
        add_rect(sl, x, Inches(1.9), Inches(4.0), Inches(4.5),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, title, x + Inches(0.2), Inches(2.1),
                    Inches(3.6), Inches(0.55), size=18, color=C_ACCENT, bold=True)
        add_textbox(sl, big, x + Inches(0.1), Inches(2.75),
                    Inches(3.8), Inches(1.0), size=36, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, sub, x + Inches(0.1), Inches(3.85),
                    Inches(3.8), Inches(0.9), size=15, color=C_GRAY,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "→ 出店タイミングは今が最適。居酒屋×発酵×定食の空白地帯は誰もいない。",
                Inches(0.65), Inches(6.5), Inches(12.3), Inches(0.65),
                size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s04_concept(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "CONCEPT", "肉酒場 然のコンセプト", 4)

    # 左：コンセプト説明
    add_textbox(sl, "業態定義", Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=14, bold=True, color=C_ACCENT)
    add_textbox(sl,
                "昼：糀漬け肉定食の専門店\n夜：発酵×肉の創作酒場",
                Inches(0.65), Inches(2.35), Inches(5.8), Inches(1.0),
                size=22, bold=True, color=C_WHITE)

    add_textbox(sl, "三つの独自性", Inches(0.65), Inches(3.5), Inches(5.8), Inches(0.45),
                size=14, bold=True, color=C_ACCENT)

    points = [
        ("糀漬け×低温調理", "前日仕込みで品質均一化。FC展開時の再現性を担保。"),
        ("昼夜二毛作モデル", "ランチ定食で固定売上 ＋ ディナー酒場でドリンク収益。"),
        ("石板焼き体験", "煙・音・ロゼ色の断面がSNS映え。来店動機を設計する。"),
    ]
    for i, (head, body) in enumerate(points):
        y = Inches(4.05 + i * 0.98)
        add_rect(sl, Inches(0.65), y + Inches(0.05), Inches(0.07), Inches(0.68), C_ACCENT)
        add_textbox(sl, head, Inches(0.88), y, Inches(2.2), Inches(0.38),
                    size=16, bold=True, color=C_WHITE)
        add_textbox(sl, body, Inches(0.88), y + Inches(0.4), Inches(5.2), Inches(0.45),
                    size=14, color=C_GRAY)

    # 右：ポジショニングマップ（matplotlib画像）
    add_textbox(sl, "ポジショニングマップ",
                Inches(7.0), Inches(1.82), Inches(6.0), Inches(0.45),
                size=14, bold=True, color=C_ACCENT)
    pos_img = make_positioning_map()
    sl.shapes.add_picture(pos_img, Inches(7.0), Inches(2.3), Inches(6.1), Inches(4.8))
    add_textbox(sl, "「健康・発酵 × 肉 × ランチ＆ディナー」は競合ゼロの完全空白地帯",
                Inches(7.0), Inches(7.12), Inches(6.1), Inches(0.35),
                size=13, bold=True, color=C_ACCENT)

    return sl


def s05_products(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "PRODUCT", "看板3品　――　これを食べに来る", 5)

    products = [
        {
            "no": "01",
            "name": "糀漬けサガリの\n石板焼き定食",
            "price": "¥1,980",
            "desc": "北海道産サガリを糀で48時間熟成。\n石板で焼き上げる煙と音が体験価値に。",
            "cost": "原価率 34.5%",
            "hook": "SNS映え No.1",
        },
        {
            "no": "02",
            "name": "糀漬け豚の\n厚切りカツ定食",
            "price": "¥1,780",
            "desc": "糀漬け豚ロースを低温調理後に揚げる。\nロゼ色の断面が口コミの起点。",
            "cost": "原価率 32.1%",
            "hook": "定食の王道×発酵",
        },
        {
            "no": "03",
            "name": "糀漬け鶏の\nロースト定食",
            "price": "¥1,680",
            "desc": "鶏もも肉を糀×ハーブで低温ロースト。\n鶏の旨みが最大限に引き出される。",
            "cost": "原価率 35.6%",
            "hook": "ヘルシー需要取込み",
        },
    ]

    for i, p in enumerate(products):
        x = Inches(0.55 + i * 4.25)
        add_rect(sl, x, Inches(1.85), Inches(4.05), Inches(5.25),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, p["no"], x + Inches(0.2), Inches(2.0),
                    Inches(0.65), Inches(0.45), size=13, bold=True, color=C_ACCENT)
        add_rect(sl, x + Inches(1.9), Inches(2.02), Inches(1.95), Inches(0.36),
                 C_ACCENT)
        add_textbox(sl, p["hook"], x + Inches(1.9), Inches(2.02),
                    Inches(1.95), Inches(0.36), size=12, bold=True,
                    color=C_BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, p["name"], x + Inches(0.2), Inches(2.5),
                    Inches(3.65), Inches(0.95), size=20, bold=True, color=C_WHITE)
        add_textbox(sl, p["price"], x + Inches(0.2), Inches(3.55),
                    Inches(3.65), Inches(0.65), size=32, bold=True, color=C_ACCENT)
        add_textbox(sl, p["desc"], x + Inches(0.2), Inches(4.3),
                    Inches(3.65), Inches(0.95), size=14, color=C_LIGHT)
        add_textbox(sl, p["cost"], x + Inches(0.2), Inches(5.35),
                    Inches(3.65), Inches(0.4), size=13, color=C_GRAY)

    add_textbox(sl,
                "全3品にライス・味噌汁・小鉢2品・漬物が付く。「体にいい定食」を徹底追求。",
                Inches(0.55), Inches(7.1), Inches(12.3), Inches(0.3),
                size=13, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return sl


def s06_competitor(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "COMPETITION", "競合5店　繁盛の法則と然の優位性", 6)

    headers = ["店舗", "月商", "ランチ", "発酵", "肉特化", "FC展開"]
    col_w   = [2.7, 1.6, 1.3, 1.3, 1.3, 1.3]
    col_x   = [0.65]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    y_h = Inches(1.82)
    add_rect(sl, Inches(0.55), y_h, Inches(12.3), Inches(0.42), C_ACCENT)
    for i, h in enumerate(headers):
        add_textbox(sl, h, Inches(col_x[i] + 0.06), y_h + Inches(0.05),
                    Inches(col_w[i]), Inches(0.32),
                    size=13, bold=True, color=C_BG)

    rows = [
        ["なかよし（恵比寿）",  "非公開", "◎", "△", "✗", "◎"],
        ["つむぎ堂（新宿）",    "1,500万", "✗", "✗", "△", "◎"],
        ["びゃく（学芸大学）",  "非公開", "✗", "◎", "✗", "△"],
        ["ひまり堂（恵比寿）",  "850万", "✗", "✗", "△", "◎"],
        ["ひまり商店（新橋）",  "1,700万", "✗", "✗", "✗", "◎"],
        ["★ 肉酒場 然",         "720万目標", "◎", "◎", "◎", "◎"],
    ]

    for r, row in enumerate(rows):
        y_r = Inches(2.26 + r * 0.68)
        bg_c = RGBColor(0x2A, 0x1E, 0x12) if r < 5 else RGBColor(0x3A, 0x2A, 0x08)
        add_rect(sl, Inches(0.55), y_r, Inches(12.3), Inches(0.64), bg_c)
        for c, cell in enumerate(row):
            color = C_ACCENT if r == 5 else (C_WHITE if c == 0 else C_LIGHT)
            add_textbox(sl, cell,
                        Inches(col_x[c] + 0.06), y_r + Inches(0.13),
                        Inches(col_w[c]), Inches(0.42),
                        size=14, bold=(r == 5), color=color)

    add_textbox(sl,
                "「健康・発酵 × 肉 × ランチ×ディナー × FC」を同時に満たすのは然だけ",
                Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38),
                size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s07_profile(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "FOUNDER", "代表者プロフィール", 7)

    add_textbox(sl, "経歴", Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    careers = [
        ("前職", "都内約25店舗展開　焼肉レストランブランド　運営本部長"),
        ("現 職", "やきとり家すみれ　運営本部長（SV体制構築・販促・業績管理）"),
        ("兼 業", "EC通販事業 経営メンバー（販売戦略・商品企画・数値分析）"),
    ]
    for i, (label, text) in enumerate(careers):
        y = Inches(2.38 + i * 0.72)
        add_rect(sl, Inches(0.65), y + Inches(0.08), Inches(0.07), Inches(0.48), C_ACCENT)
        add_textbox(sl, label, Inches(0.87), y, Inches(1.0), Inches(0.38),
                    size=13, color=C_ACCENT, bold=True)
        add_textbox(sl, text, Inches(1.98), y, Inches(4.4), Inches(0.6),
                    size=15, color=C_WHITE)

    add_textbox(sl, "5つの強み", Inches(7.2), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    strengths = [
        "複数店舗の運営経験（本部長として売上・利益・人材を総合管理）",
        "現場と経営の両立（現場課題と経営数値の両方を深く理解）",
        "数字に基づく改善力（客数・客単価・原価・人件費を分解改善）",
        "商品企画・販促経験（SNS・LINE・キャンペーンによる集客実績）",
        "人材育成と仕組み化（属人化しない標準化・教育体制の構築）",
    ]
    for i, s in enumerate(strengths):
        y = Inches(2.38 + i * 0.75)
        add_rect(sl, Inches(7.1), y + Inches(0.08), Inches(5.9), Inches(0.58),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, f"0{i+1}  {s}", Inches(7.25), y + Inches(0.1),
                    Inches(5.6), Inches(0.54), size=14, color=C_LIGHT)

    add_rect(sl, Inches(0.55), Inches(6.05), Inches(12.3), Inches(1.05),
             RGBColor(0x2A, 0x1E, 0x12))
    add_textbox(sl,
                "「飲食の現場と数字を両方知っている人間が、\n"
                "　ゼロから作るブランドだから強い。」",
                Inches(0.85), Inches(6.12), Inches(11.7), Inches(0.9),
                size=18, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s08_pl(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "FINANCIALS", "収支計画　――　堅実な数字設計", 8)

    add_textbox(sl, "月次売上目標（安定期・30席）",
                Inches(0.65), Inches(1.82), Inches(6.0), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    sales_rows = [
        ("ランチ", "2,000円", "30席", "1.8回転", "25日", "270万円"),
        ("ディナー", "5,000円", "30席", "1.2回転", "25日", "450万円"),
        ("合計", "", "", "", "", "720万円"),
    ]
    s_headers = ["", "客単価", "席数", "回転", "営業日", "売上"]
    s_col_w = [1.15, 1.05, 0.85, 0.95, 0.85, 1.05]
    s_col_x = [0.65]
    for w in s_col_w[:-1]:
        s_col_x.append(s_col_x[-1] + w)

    add_rect(sl, Inches(0.55), Inches(2.3), Inches(6.1), Inches(0.36), C_ACCENT)
    for i, h in enumerate(s_headers):
        add_textbox(sl, h, Inches(s_col_x[i]+0.04), Inches(2.32),
                    Inches(s_col_w[i]), Inches(0.3), size=12, bold=True, color=C_BG)

    for r, row in enumerate(sales_rows):
        y = Inches(2.68 + r * 0.52)
        bg_c = RGBColor(0x3A, 0x2A, 0x08) if r == 2 else RGBColor(0x2A, 0x1E, 0x12)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.5), bg_c)
        for c, cell in enumerate(row):
            col = C_ACCENT if (r == 2 and c == 5) else C_WHITE
            add_textbox(sl, cell, Inches(s_col_x[c]+0.04), y + Inches(0.08),
                        Inches(s_col_w[c]), Inches(0.36), size=13,
                        bold=(r == 2), color=col)

    add_textbox(sl, "フェーズ別収支",
                Inches(0.65), Inches(4.2), Inches(6.0), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    phases = [
        ("立ち上げ期", "0〜3ヶ月", "360万", "▲72万", "C_RED"),
        ("成長期",     "3〜6ヶ月", "504万", "＋14万", "C_GRAY"),
        ("安定期",     "6ヶ月〜",  "648万", "＋142万", "C_GREEN"),
    ]
    for i, (ph, period, sales, profit, cc) in enumerate(phases):
        y = Inches(4.72 + i * 0.65)
        pcolor = C_RED if cc == "C_RED" else (C_GREEN if cc == "C_GREEN" else C_GRAY)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.6),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, ph, Inches(0.7), y + Inches(0.1),
                    Inches(1.6), Inches(0.42), size=15, bold=True, color=C_WHITE)
        add_textbox(sl, period, Inches(2.2), y + Inches(0.1),
                    Inches(1.3), Inches(0.42), size=14, color=C_GRAY)
        add_textbox(sl, sales, Inches(3.4), y + Inches(0.1),
                    Inches(1.3), Inches(0.42), size=14, color=C_LIGHT)
        add_textbox(sl, profit, Inches(4.7), y + Inches(0.1),
                    Inches(1.5), Inches(0.42), size=16, bold=True, color=pcolor)

    # 右：損益分岐
    add_rect(sl, Inches(7.1), Inches(1.82), Inches(6.0), Inches(5.3),
             RGBColor(0x2A, 0x1E, 0x12))
    add_textbox(sl, "損益分岐点分析",
                Inches(7.25), Inches(1.96), Inches(5.7), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    bep_items = [
        ("損益分岐点売上", "460万円/月"),
        ("目標売上", "720万円/月"),
        ("目標対比", "64%で黒字転換"),
        ("固定費合計", "約175万円/月"),
        ("FL比率（安定期）", "48.6%"),
        ("安定期営業利益", "約142万円/月"),
        ("月次返済（友人＋公庫）", "約20万円/月"),
        ("返済後手残り（安定期）", "約122万円/月"),
    ]
    for i, (k, v) in enumerate(bep_items):
        y = Inches(2.5 + i * 0.56)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(0.52),
                 RGBColor(0x22, 0x18, 0x0E) if i % 2 == 0 else RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, k, Inches(7.25), y + Inches(0.08),
                    Inches(3.3), Inches(0.38), size=13, color=C_GRAY)
        vc = C_ACCENT if "利益" in k or "手残り" in k else C_WHITE
        add_textbox(sl, v, Inches(10.4), y + Inches(0.08),
                    Inches(2.5), Inches(0.38), size=14, bold=True,
                    color=vc, align=PP_ALIGN.RIGHT)

    return sl


def s09_finance(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "FUNDING", "資金調達計画", 9)

    funding = [
        ("代表 自己資金", "630万円", "90%持分（株式）"),
        ("友人 出資", "70万円", "10%持分（株式）"),
        ("友人 株主借入金", "1,430万円", "金利1%・10年返済"),
        ("日本政策金融公庫", "500万円", "金利2%・7年返済"),
    ]
    add_textbox(sl, "調達構成",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    for i, (name, amount, note) in enumerate(funding):
        y = Inches(2.38 + i * 0.75)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.68),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, name, Inches(0.7), y + Inches(0.08),
                    Inches(2.3), Inches(0.55), size=15, color=C_LIGHT)
        add_textbox(sl, amount, Inches(2.95), y + Inches(0.06),
                    Inches(1.6), Inches(0.58), size=18, bold=True, color=C_ACCENT)
        add_textbox(sl, note, Inches(4.45), y + Inches(0.1),
                    Inches(2.1), Inches(0.52), size=13, color=C_GRAY)

    add_rect(sl, Inches(0.55), Inches(5.4), Inches(6.1), Inches(0.72),
             RGBColor(0x3A, 0x2A, 0x08))
    add_textbox(sl, "調達総額", Inches(0.7), Inches(5.5),
                Inches(2.1), Inches(0.55), size=16, bold=True, color=C_WHITE)
    add_textbox(sl, "2,630万円", Inches(2.95), Inches(5.46),
                Inches(2.8), Inches(0.62), size=22, bold=True, color=C_ACCENT)

    add_textbox(sl, "資金使途",
                Inches(7.2), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    uses = [
        ("物件取得費（保証金・礼金）", "350万円"),
        ("内装・設備工事", "870万円"),
        ("厨房機器・備品", "200万円"),
        ("仕入・食材初回発注", "50万円"),
        ("採用・オープン前研修", "80万円"),
        ("マーケティング（3ヶ月分）", "45万円"),
        ("運転資金（3ヶ月分）", "405万円"),
        ("手元残金（緊急予備）", "630万円"),
    ]
    for i, (use, amt) in enumerate(uses):
        y = Inches(2.38 + i * 0.58)
        bg_c = RGBColor(0x3A, 0x2A, 0x08) if "手元" in use else RGBColor(0x2A, 0x1E, 0x12)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(0.54), bg_c)
        col = C_ACCENT if "手元" in use else C_LIGHT
        add_textbox(sl, use, Inches(7.25), y + Inches(0.08),
                    Inches(3.9), Inches(0.4), size=14, color=col)
        add_textbox(sl, amt, Inches(11.1), y + Inches(0.08),
                    Inches(1.8), Inches(0.4), size=15, bold=True,
                    color=col, align=PP_ALIGN.RIGHT)

    add_textbox(sl,
                "手元残金630万円は立ち上げ期の赤字（▲216万）を十分カバー。財務安全性は高い。",
                Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38),
                size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s10_roadmap(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "ROADMAP", "展開ロードマップ　7年で100店舗へ", 10)

    phases = [
        {
            "phase": "実証期",
            "period": "1〜2年目",
            "stores": "1〜2店舗",
            "action": "直営1号店で型をつくる\nQSC・レシピ・教育マニュアル完成\nインフルエンサー戦略で認知獲得",
            "kpi": "月商黒字化\nリピート率60%以上",
        },
        {
            "phase": "横展開期",
            "period": "3年目",
            "stores": "2〜5店舗",
            "action": "直営2〜3号店 ＋ FC募集開始\nFC加盟オーナー教育体制構築\n三軒茶屋・目黒・自由が丘エリアへ",
            "kpi": "FC1号店開業\nグループ月商2,000万超",
        },
        {
            "phase": "加速期",
            "period": "4〜5年目",
            "stores": "20〜50店舗",
            "action": "FC中心に年15〜20店舗ペース\n魚酒場 然・鶏酒場 然も展開開始\nSV体制・本部機能を整備",
            "kpi": "FC比率70%以上\n年商10億超",
        },
        {
            "phase": "確立期",
            "period": "6〜7年目",
            "stores": "100店舗",
            "action": "全国展開・上場検討\nブランド価値の確立\n「然」を食文化のアイコンへ",
            "kpi": "100店舗達成\n上場準備",
        },
    ]

    for i, p in enumerate(phases):
        x = Inches(0.55 + i * 3.2)
        add_rect(sl, x, Inches(1.82), Inches(3.05), Inches(0.48), C_ACCENT)
        add_textbox(sl, f"{p['phase']}　{p['period']}",
                    x + Inches(0.1), Inches(1.85), Inches(2.85), Inches(0.42),
                    size=14, bold=True, color=C_BG)
        add_rect(sl, x, Inches(2.3), Inches(3.05), Inches(5.1),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, p["stores"],
                    x + Inches(0.12), Inches(2.4), Inches(2.85), Inches(0.65),
                    size=28, bold=True, color=C_ACCENT if i == 3 else C_WHITE)
        add_line(sl, x + Inches(0.12), Inches(3.12), Inches(2.78), color=C_ACCENT,
                 height=Pt(1.0))
        add_textbox(sl, p["action"],
                    x + Inches(0.12), Inches(3.25), Inches(2.85), Inches(2.2),
                    size=13, color=C_LIGHT)
        add_rect(sl, x, Inches(5.4), Inches(3.05), Inches(0.85),
                 RGBColor(0x1A, 0x12, 0x05))
        add_textbox(sl, "KPI: " + p["kpi"],
                    x + Inches(0.12), Inches(5.45), Inches(2.85), Inches(0.72),
                    size=13, color=C_ACCENT)

    for i in range(3):
        add_textbox(sl, "→", Inches(3.38 + i * 3.2), Inches(3.85),
                    Inches(0.28), Inches(0.42), size=22, color=C_ACCENT,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl, "FC展開により資本効率を最大化。オーナーのノウハウを型にして横展開する。",
                Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38),
                size=14, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return sl


def s11_strategy(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "MARKETING", "集客・マーケティング戦略", 11)

    strategies = [
        {
            "title": "① インフルエンサー戦略",
            "ref": "ひまり商店モデル（初月1,700万円）",
            "body": "開業3ヶ月前からマイクロ〜メガ\nインフルエンサーをリスト化。\n無料招待→動画・リール投稿で\nオープン前から認知を爆発させる。",
        },
        {
            "title": "② Instagram中心のSNS設計",
            "ref": "石板の煙・音・断面美を動画化",
            "body": "撮られることを前提にした\n商品・空間・BGMの設計。\n「映え」ではなく「体験動画」\nとして拡散される仕掛けを作る。",
        },
        {
            "title": "③ MEO（Googleマップ）対策",
            "ref": "「肉酒場 三軒茶屋」で1位獲得",
            "body": "業態説明に「糀漬け肉定食居酒屋」\nを明記。写真・口コミ・投稿頻度\nで上位表示を維持する。\nAI検索（AIO）にも対応した記述。",
        },
        {
            "title": "④ ランチ×ディナー来店設計",
            "ref": "一人の客が昼も夜も来る店に",
            "body": "ランチで日常客を獲得し、\n夜の酒場に引き上げる。\nLINE公式・スタンプカードで\nリピートを仕組み化する。",
        },
    ]

    for i, s in enumerate(strategies):
        col = i % 2
        row = i // 2
        x = Inches(0.55 + col * 6.45)
        y = Inches(1.82 + row * 2.7)
        add_rect(sl, x, y, Inches(6.2), Inches(2.55), RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, s["title"], x + Inches(0.18), y + Inches(0.12),
                    Inches(5.85), Inches(0.45), size=17, bold=True, color=C_WHITE)
        add_textbox(sl, s["ref"], x + Inches(0.18), y + Inches(0.6),
                    Inches(5.85), Inches(0.32), size=13, color=C_ACCENT, italic=True)
        add_line(sl, x + Inches(0.18), y + Inches(0.97), Inches(5.8),
                 color=C_GRAY, height=Pt(0.5))
        add_textbox(sl, s["body"], x + Inches(0.18), y + Inches(1.08),
                    Inches(5.85), Inches(1.35), size=14, color=C_LIGHT)

    return sl


def s13_open_strategy(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "LAUNCH STRATEGY", "オープン戦略　――　開業前から火をつける", 13)

    timeline = [
        ("-6ヶ月", "仕込み期", [
            "物件契約・内装着工",
            "Instagramアカウント開設（コンセプト発信）",
            "インフルエンサーリスト作成（50名以上）",
            "メニュー最終決定・レシピ標準化",
        ]),
        ("-3ヶ月", "認知獲得期", [
            "マイクロインフルエンサー（1万〜10万）20名招待",
            "工事中の様子・仕込み動画をリール投稿",
            "Googleビジネスプロフィール登録・MEO開始",
            "プレスリリース・食メディアへのアプローチ",
        ]),
        ("-1ヶ月", "加熱期", [
            "メガインフルエンサー（10万〜）5名招待",
            "プレオープン（関係者・常連候補・地元住民）",
            "Instagramカウントダウン投稿",
            "Googleマップ写真・口コミ初期投稿の準備",
        ]),
        ("OPEN", "爆発期", [
            "オープン当日ライブ配信",
            "初月は意図的に「並ばせる」演出で話題化",
            "日次で売上・客数・SNS反応をモニタリング",
            "口コミ返信を毎日・MEO評価を週次チェック",
        ]),
    ]

    for i, (time, phase, items) in enumerate(timeline):
        x = Inches(0.45 + i * 3.22)
        bg_c = C_ACCENT if time == "OPEN" else RGBColor(0x3A, 0x2A, 0x08)
        txt_c = C_BG if time == "OPEN" else C_ACCENT
        add_rect(sl, x, Inches(1.82), Inches(3.05), Inches(0.44), bg_c)
        add_textbox(sl, time, x + Inches(0.06), Inches(1.85),
                    Inches(2.93), Inches(0.38), size=16, bold=True,
                    color=txt_c, align=PP_ALIGN.CENTER)
        add_textbox(sl, phase, x + Inches(0.12), Inches(2.32),
                    Inches(2.85), Inches(0.4), size=15, bold=True, color=C_WHITE)
        add_rect(sl, x, Inches(2.78), Inches(3.05), Inches(4.55),
                 RGBColor(0x2A, 0x1E, 0x12))
        for j, item in enumerate(items):
            add_rect(sl, x + Inches(0.14), Inches(2.9 + j * 1.05),
                     Inches(0.07), Inches(0.74), C_ACCENT)
            add_textbox(sl, item, x + Inches(0.3), Inches(2.9 + j * 1.05),
                        Inches(2.6), Inches(0.9), size=13, color=C_LIGHT)

    add_textbox(sl,
                "参考：ひまり商店（新橋）は開業前からのインフルエンサー戦略で初月1,700万円を達成",
                Inches(0.55), Inches(7.1), Inches(12.3), Inches(0.3),
                size=12, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return sl


def s14_fc_strategy(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "FC STRATEGY", "FC展開戦略　――　型をつくり、型で増やす", 14)

    add_textbox(sl, "FCパッケージ設計（予定）",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    fc_items = [
        ("加盟金", "200万円", "ブランド使用権・初期研修込み"),
        ("ロイヤリティ", "売上の3%", "本部サポート・仕入支援含む"),
        ("保証金", "100万円", "退会時返還"),
        ("加盟者の初期投資目安", "1,500〜2,000万円", "物件・内装・設備込み"),
        ("FC加盟者の想定月商", "600〜720万円", "直営モデルの再現"),
        ("FCオーナーの想定手残り", "月60〜90万円", "ロイヤリティ・返済後"),
    ]
    for i, (k, v, note) in enumerate(fc_items):
        y = Inches(2.38 + i * 0.67)
        bg_c = RGBColor(0x3A, 0x2A, 0x08) if "手残り" in k else RGBColor(0x2A, 0x1E, 0x12)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.62), bg_c)
        add_textbox(sl, k, Inches(0.7), y + Inches(0.1),
                    Inches(2.6), Inches(0.44), size=13, color=C_GRAY)
        vc = C_ACCENT if "手残り" in k else C_WHITE
        add_textbox(sl, v, Inches(3.15), y + Inches(0.06),
                    Inches(1.85), Inches(0.52), size=16, bold=True, color=vc)
        add_textbox(sl, note, Inches(4.85), y + Inches(0.1),
                    Inches(1.7), Inches(0.44), size=11, color=C_GRAY)

    add_textbox(sl, "FC加盟者ターゲット像",
                Inches(7.3), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    targets = [
        ("飲食経験者の独立志望者", "元料理人・飲食チェーン出身者。\n現場力があり、ブランドの型を活かせる。"),
        ("副業・投資目的の個人オーナー", "他業種からの参入。オーナー不在でも\n回るオペレーションを本部が支援。"),
        ("地方への展開パートナー", "東京モデルを地方に持ち込む\nエリアFCオーナー。複数店舗可。"),
    ]
    for i, (title, desc) in enumerate(targets):
        y = Inches(2.38 + i * 1.52)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(1.38),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_rect(sl, Inches(7.1), y, Inches(0.09), Inches(1.38), C_ACCENT)
        add_textbox(sl, title, Inches(7.32), y + Inches(0.1),
                    Inches(5.6), Inches(0.42), size=16, bold=True, color=C_WHITE)
        add_textbox(sl, desc, Inches(7.32), y + Inches(0.57),
                    Inches(5.6), Inches(0.72), size=14, color=C_LIGHT)

    add_textbox(sl, "FC本部収益（100店舗時）",
                Inches(7.3), Inches(6.88), Inches(5.8), Inches(0.3),
                size=13, bold=True, color=C_ACCENT)
    add_textbox(sl,
                "100店×720万×3%ロイヤリティ ＝ 月2,160万円 / 年2.6億円",
                Inches(7.1), Inches(7.15), Inches(6.0), Inches(0.3),
                size=13, color=C_WHITE)

    return sl


def s15_location(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "LOCATION STRATEGY", "立地・出店戦略　――　勝てる場所を選ぶ", 15)

    add_textbox(sl, "物件選定スコアリング基準",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    criteria = [
        ("乗降者数", "5万人/日以上", "◎"),
        ("駅徒歩", "5分以内・路面優先", "◎"),
        ("坪数", "15〜25坪（理想20坪）", "◎"),
        ("家賃", "50万円以内（坪2.5万以下）", "◎"),
        ("設備", "重飲食対応・居抜き優先", "◎"),
        ("競合", "同業態が近くに集中していない", "△"),
    ]
    for i, (k, v, mark) in enumerate(criteria):
        y = Inches(2.38 + i * 0.65)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.6),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, k, Inches(0.7), y + Inches(0.1),
                    Inches(1.4), Inches(0.42), size=14, color=C_GRAY)
        add_textbox(sl, v, Inches(2.0), y + Inches(0.1),
                    Inches(3.3), Inches(0.42), size=14, color=C_WHITE)
        mc = C_ACCENT if mark == "◎" else C_GRAY
        add_textbox(sl, mark, Inches(5.8), y + Inches(0.08),
                    Inches(0.65), Inches(0.44), size=18, bold=True, color=mc,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl, "優先出店エリア（1号店候補）",
                Inches(7.3), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    areas = [
        ("★ 三軒茶屋", "約10万人/日", "住宅街×夜の人通り◎\n飲食激戦区だが差別化余地大\n渋谷から2駅の好立地"),
        ("★ 目黒",     "約12万人/日", "30〜40代富裕層多い\n平均客単価が高く取りやすい\nランチ需要も強い"),
        ("★ 自由が丘", "約8万人/日",  "感度高い客層・SNS拡散◎\nカフェ・グルメ激戦区\nブランドイメージ向上に寄与"),
    ]
    for i, (name, riders, desc) in enumerate(areas):
        y = Inches(2.38 + i * 1.62)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(1.48),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, name, Inches(7.25), y + Inches(0.1),
                    Inches(2.3), Inches(0.44), size=18, bold=True, color=C_ACCENT)
        add_textbox(sl, riders, Inches(9.5), y + Inches(0.12),
                    Inches(2.4), Inches(0.4), size=15, color=C_WHITE,
                    align=PP_ALIGN.RIGHT)
        add_textbox(sl, desc, Inches(7.25), y + Inches(0.6),
                    Inches(5.7), Inches(0.8), size=14, color=C_LIGHT)

    add_textbox(sl,
                "1号店は三軒茶屋〜目黒エリアで検討。居抜き物件を優先し初期投資を抑制する。",
                Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38),
                size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s16_brand_sns(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "BRAND & SNS", "ブランド・SNS戦略　――　「然」を育てる", 16)

    add_textbox(sl, "ブランドアーキテクチャ",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    add_rect(sl, Inches(2.1), Inches(2.38), Inches(2.95), Inches(0.56), C_ACCENT)
    add_textbox(sl, "然（ぜん）", Inches(2.1), Inches(2.41),
                Inches(2.95), Inches(0.48), size=20, bold=True,
                color=C_BG, align=PP_ALIGN.CENTER)

    branches = [
        ("肉酒場 然", "1号店〜\n糀漬け肉×石板焼き", Inches(0.4)),
        ("魚酒場 然", "2年目〜\n魚介×発酵", Inches(2.6)),
        ("鶏酒場 然", "3年目〜\nやきとり×糀", Inches(4.8)),
    ]
    for name, desc, x in branches:
        add_rect(sl, x, Inches(3.55), Inches(2.05), Inches(0.46),
                 RGBColor(0x3A, 0x2A, 0x08))
        add_textbox(sl, name, x + Inches(0.06), Inches(3.58),
                    Inches(1.93), Inches(0.4), size=14, bold=True,
                    color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(sl, desc, x + Inches(0.06), Inches(4.08),
                    Inches(1.93), Inches(0.58), size=12, color=C_GRAY,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl, "SEO / MEO / AIO設計",
                Inches(0.65), Inches(4.82), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)
    seo_items = [
        ("SEO", "「肉酒場」「糀漬け定食」＋エリア名でコンテンツ設計"),
        ("MEO", "Googleマップ説明に「糀漬け肉定食居酒屋」を明記。写真・口コミを週次更新"),
        ("AIO", "AI検索で「糀漬け肉定食の居酒屋」として説明されやすい構造的な記述"),
        ("指名", "「然（ぜん）」を固有名詞として育て、リピーター指名検索を増やす"),
    ]
    for i, (tag, desc) in enumerate(seo_items):
        y = Inches(5.35 + i * 0.48)
        add_rect(sl, Inches(0.55), y, Inches(0.58), Inches(0.4), C_ACCENT)
        add_textbox(sl, tag, Inches(0.55), y + Inches(0.02),
                    Inches(0.58), Inches(0.36), size=11, bold=True,
                    color=C_BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, desc, Inches(1.22), y,
                    Inches(5.3), Inches(0.42), size=13, color=C_LIGHT)

    add_textbox(sl, "SNSコンテンツ設計（Instagram）",
                Inches(7.3), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    contents = [
        ("体験動画（リール）", "石板の煙・音・肉の断面。\n「見ていて気持ちいい」が拡散の鍵。", "週3本"),
        ("仕込みの裏側", "糀漬けの工程・仕込み動画。\n「作り手が見える」信頼感を演出。", "週2本"),
        ("インフルエンサーUGC", "招待者の投稿をリポスト。\n他者目線の口コミが最強の広告。", "随時"),
        ("オーナーの想い発信", "開業ストーリー・日々の学び。\n「この人から買いたい」を育てる。", "週1本"),
    ]
    for i, (title, body, freq) in enumerate(contents):
        y = Inches(2.38 + i * 1.24)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(1.14),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_textbox(sl, title, Inches(7.25), y + Inches(0.07),
                    Inches(3.9), Inches(0.4), size=15, bold=True, color=C_WHITE)
        add_rect(sl, Inches(11.4), y + Inches(0.1), Inches(1.4), Inches(0.32),
                 RGBColor(0x3A, 0x2A, 0x08))
        add_textbox(sl, freq, Inches(11.4), y + Inches(0.1),
                    Inches(1.4), Inches(0.32), size=11, color=C_ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, body, Inches(7.25), y + Inches(0.5),
                    Inches(5.7), Inches(0.57), size=13, color=C_LIGHT)

    return sl


def s17_operations(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "OPERATIONS", "オペレーション設計　――　型が利益を生む", 17)

    add_textbox(sl, "仕込みフロー（前日→当日）",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    steps = [
        ("前日 15:00", "サガリ・豚・鶏を糀に漬け込み（各48h）"),
        ("前日 17:00", "低温調理器にセット（サガリ55℃×4h）"),
        ("当日 10:00", "仕込み完了・ポーション確認・ライン設定"),
        ("11:30", "ランチ営業開始（石板+トッピング仕上げ30秒）"),
        ("14:00", "ランチ終了・ディナー仕込み補充・清掃"),
        ("17:30", "ディナー営業開始"),
        ("22:30", "閉店・翌日仕込み開始"),
    ]
    for i, (time, action) in enumerate(steps):
        y = Inches(2.38 + i * 0.6)
        add_rect(sl, Inches(0.55), y, Inches(6.1), Inches(0.56),
                 RGBColor(0x2A, 0x1E, 0x12) if i % 2 == 0 else RGBColor(0x22, 0x18, 0x0E))
        add_textbox(sl, time, Inches(0.7), y + Inches(0.1),
                    Inches(1.35), Inches(0.4), size=13, bold=True, color=C_ACCENT)
        add_textbox(sl, action, Inches(2.1), y + Inches(0.1),
                    Inches(4.4), Inches(0.4), size=14, color=C_WHITE)

    add_textbox(sl, "FC展開を可能にする3つの仕組み",
                Inches(7.3), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    pillars = [
        ("① 低温調理×前日仕込み",
         "温度・時間が数値化されているため\n誰が作っても同じ品質に仕上がる。\n料理経験が浅いスタッフでも対応可能。"),
        ("② 石板仕上げ（30秒）",
         "調理の最終工程は石板で30秒焼くだけ。\nオペレーションが極めてシンプルで\nホール兼務も可能なスキーム。"),
        ("③ レシピ・マニュアルの完全数値化",
         "仕込み温度・時間・グラム数を\n全てマニュアル化。SVが巡回研修で\n品質を均一化する体制を構築。"),
    ]
    for i, (title, body) in enumerate(pillars):
        y = Inches(2.38 + i * 1.62)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(1.48),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_rect(sl, Inches(7.1), y, Inches(0.09), Inches(1.48), C_ACCENT)
        add_textbox(sl, title, Inches(7.32), y + Inches(0.1),
                    Inches(5.6), Inches(0.44), size=16, bold=True, color=C_WHITE)
        add_textbox(sl, body, Inches(7.32), y + Inches(0.6),
                    Inches(5.6), Inches(0.8), size=14, color=C_LIGHT)

    add_textbox(sl,
                "「低温調理×石板仕上げ」は職人不要のオペレーションを実現し、FC展開の最大の武器になる。",
                Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38),
                size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s18_data_kpi(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "DATA MANAGEMENT", "データ経営　――　感覚ではなく数字で動く", 18)

    add_textbox(sl, "KPI管理サイクル",
                Inches(0.65), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    cycles = [
        ("日次", [
            "客数・客単価・売上",
            "時間帯別回転数",
            "メニュー別注文数",
            "SNSエンゲージメント",
        ]),
        ("週次", [
            "原価率・人件費率",
            "FL比率の確認",
            "口コミ・MEO評価",
            "前週比・前年比",
        ]),
        ("月次", [
            "損益計算書の確認",
            "リピート率・新規比率",
            "メニュー見直し判断",
            "広告ROI測定",
        ]),
    ]
    for i, (freq, items) in enumerate(cycles):
        x = Inches(0.55 + i * 2.25)
        add_rect(sl, x, Inches(2.38), Inches(2.05), Inches(0.44), C_ACCENT)
        add_textbox(sl, freq, x, Inches(2.41), Inches(2.05), Inches(0.38),
                    size=15, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_rect(sl, x, Inches(2.82), Inches(2.05), Inches(3.15),
                 RGBColor(0x2A, 0x1E, 0x12))
        for j, item in enumerate(items):
            add_textbox(sl, f"• {item}", x + Inches(0.12), Inches(2.96 + j * 0.7),
                        Inches(1.85), Inches(0.62), size=13, color=C_LIGHT)

    add_textbox(sl, "重要KPI目標値（1号店・安定期）",
                Inches(7.3), Inches(1.85), Inches(5.8), Inches(0.45),
                size=15, bold=True, color=C_ACCENT)

    kpis = [
        ("月商目標", "720万円", "安定期"),
        ("損益分岐点", "460万円", "目標の64%"),
        ("ランチ回転数", "1.8回転", "30席×25日"),
        ("ディナー回転数", "1.2回転", "30席×25日"),
        ("原価率", "30%以下", "FL管理"),
        ("人件費率", "18〜22%", "体制次第"),
        ("FL比率", "48〜52%", "目標50%以下"),
        ("リピート率", "60%以上", "6ヶ月後目標"),
        ("Googleマップ評価", "4.0以上維持", "MEO管理"),
        ("Instagram フォロワー", "1万人（1年目）", "SNS資産"),
    ]
    for i, (k, v, note) in enumerate(kpis):
        y = Inches(2.38 + i * 0.49)
        bg_c = RGBColor(0x3A, 0x2A, 0x08) if i % 2 == 0 else RGBColor(0x2A, 0x1E, 0x12)
        add_rect(sl, Inches(7.1), y, Inches(6.0), Inches(0.46), bg_c)
        add_textbox(sl, k, Inches(7.25), y + Inches(0.05),
                    Inches(2.6), Inches(0.38), size=13, color=C_GRAY)
        add_textbox(sl, v, Inches(9.75), y + Inches(0.04),
                    Inches(1.8), Inches(0.4), size=14, bold=True, color=C_WHITE)
        add_textbox(sl, note, Inches(11.45), y + Inches(0.07),
                    Inches(1.45), Inches(0.35), size=11, color=C_GRAY,
                    align=PP_ALIGN.RIGHT)

    add_textbox(sl,
                "参考：つむぎ堂は毎日注文数・売上をスタッフと共有するデータ経営で坪月商57万円を実現",
                Inches(0.55), Inches(7.1), Inches(12.3), Inches(0.3),
                size=12, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return sl


def s_growth_chart(prs, slide_num):
    """成長・財務チャートスライド（グラフ2本: 年商＋営業利益）"""
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "10-YEAR GROWTH", "10年財務成長シナリオ", slide_num)

    # サマリーKPI（4ボックス）
    # 7年目の実値を計算して表示
    y7_rev    = REVENUE[6]
    y7_profit = PROFIT[6]
    y7_rate   = MARGIN_RATE[6]
    y10_rev   = REVENUE[9]
    kpis = [
        ("7年目 グループ年商",   f"{y7_rev:.0f}億円"),
        ("7年目 営業利益",       f"{y7_profit:.1f}億円"),
        ("7年目 実効利益率",     f"{y7_rate*100:.0f}%"),
        ("10年目 グループ年商",  f"{y10_rev:.0f}億円"),
    ]
    for i, (label, val) in enumerate(kpis):
        x = Inches(0.55 + i * 3.2)
        add_rect(sl, x, Inches(1.82), Inches(3.05), Inches(0.92),
                 RGBColor(0x3A, 0x2A, 0x08))
        add_textbox(sl, label, x + Inches(0.12), Inches(1.88),
                    Inches(2.85), Inches(0.32), size=12, color=C_GRAY,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, val, x + Inches(0.12), Inches(2.2),
                    Inches(2.85), Inches(0.5), size=22, bold=True,
                    color=C_ACCENT, align=PP_ALIGN.CENTER)

    # グラフ2枚を横並び（大サイズ）
    charts = [
        (make_revenue_chart, Inches(0.45),  Inches(2.9), Inches(6.2), Inches(4.45)),
        (make_profit_chart,  Inches(6.85),  Inches(2.9), Inches(6.2), Inches(4.45)),
    ]
    for fn, left, top, w, h in charts:
        img_buf = fn()
        sl.shapes.add_picture(img_buf, left, top, w, h)

    return sl


def s19_risk(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "RISK MANAGEMENT", "リスクと対策　――　想定して備える", 20)

    risks = [
        {
            "risk": "① 立ち上げ期の売上未達",
            "level": "高",
            "level_c": C_RED,
            "cause": "新店認知不足・口コミ形成に時間",
            "action": "開業3ヶ月前からインフルエンサー戦略を実行。\n手元残金630万円で最大6ヶ月の赤字に耐えられる財務設計。",
        },
        {
            "risk": "② 人材採用・定着の問題",
            "level": "中",
            "level_c": C_ACCENT,
            "cause": "飲食業界の人手不足・離職率の高さ",
            "action": "時給を相場より50〜100円高く設定し採用競争力を確保。\n仕組み化により属人依存を排除。常時採用を継続。",
        },
        {
            "risk": "③ 原材料費の高騰",
            "level": "中",
            "level_c": C_ACCENT,
            "cause": "牛サガリの価格変動・円安影響",
            "action": "複数仕入れ先を確保し価格交渉力を維持。\n鶏・豚メニューの比率調整で原価率をコントロール。",
        },
        {
            "risk": "④ 競合の模倣・追随",
            "level": "低",
            "level_c": C_GREEN,
            "cause": "糀漬け肉定食業態の競合参入",
            "action": "ブランド「然」の商標登録で保護（第43類・飲食）。\n先行者優位と顧客ロイヤリティで差別化を維持。",
        },
        {
            "risk": "⑤ FC加盟者の品質管理",
            "level": "中",
            "level_c": C_ACCENT,
            "cause": "FC展開後のQSCばらつき",
            "action": "SV定期巡回（月1回以上）・抜き打ち覆面調査の実施。\n数値基準を下回る店舗には改善プログラムを発動。",
        },
    ]

    for i, r in enumerate(risks):
        y = Inches(1.82 + i * 1.03)
        add_rect(sl, Inches(0.55), y, Inches(12.3), Inches(0.96),
                 RGBColor(0x2A, 0x1E, 0x12))
        add_rect(sl, Inches(0.55), y, Inches(0.52), Inches(0.96), r["level_c"])
        add_textbox(sl, r["level"], Inches(0.55), y + Inches(0.3),
                    Inches(0.52), Inches(0.38), size=13, bold=True,
                    color=C_BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, r["risk"], Inches(1.18), y + Inches(0.06),
                    Inches(3.1), Inches(0.42), size=15, bold=True, color=C_WHITE)
        add_textbox(sl, r["cause"], Inches(1.18), y + Inches(0.52),
                    Inches(3.1), Inches(0.38), size=13, color=C_GRAY, italic=True)
        add_rect(sl, Inches(4.45), y + Inches(0.1), Inches(0.52), Inches(0.3),
                 C_ACCENT)
        add_textbox(sl, "対策", Inches(4.45), y + Inches(0.1),
                    Inches(0.52), Inches(0.3), size=10, bold=True,
                    color=C_BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, r["action"], Inches(5.06), y + Inches(0.06),
                    Inches(7.6), Inches(0.86), size=13, color=C_LIGHT)

    return sl


def s20_action_plan(prs):
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "ACTION PLAN", "今後90日のアクションプラン", 21)

    phases_90 = [
        {
            "period": "0〜30日",
            "theme": "仕込みと準備",
            "color": RGBColor(0x3A, 0x2A, 0x08),
            "items": [
                "□ 会社設立（登記・口座開設）",
                "□ 「肉酒場 然」商標調査・出願",
                "□ 物件情報収集（三軒茶屋・目黒・自由が丘）",
                "□ 日本政策金融公庫への相談・申込",
                "□ Instagramアカウント開設",
            ],
        },
        {
            "period": "30〜60日",
            "theme": "物件決定と資金確定",
            "color": RGBColor(0x3A, 0x28, 0x05),
            "items": [
                "□ 物件内見・スコアリング評価",
                "□ 公庫面談・融資確定",
                "□ 物件契約・内装業者選定",
                "□ インフルエンサーリスト作成（50名）",
                "□ メニュー最終試作・原価計算確定",
            ],
        },
        {
            "period": "60〜90日",
            "theme": "着工と認知獲得",
            "color": RGBColor(0x3A, 0x2E, 0x08),
            "items": [
                "□ 内装・厨房工事着工",
                "□ スタッフ採用開始",
                "□ インフルエンサー招待DM開始",
                "□ Googleビジネスプロフィール登録",
                "□ POSシステム・決済端末導入",
            ],
        },
    ]

    for i, p in enumerate(phases_90):
        x = Inches(0.55 + i * 4.25)
        add_rect(sl, x, Inches(1.82), Inches(4.1), Inches(0.46), p["color"])
        add_textbox(sl, p["period"],
                    x + Inches(0.12), Inches(1.85), Inches(1.6), Inches(0.4),
                    size=14, bold=True, color=C_ACCENT)
        add_textbox(sl, p["theme"],
                    x + Inches(1.7), Inches(1.85), Inches(2.3), Inches(0.4),
                    size=13, color=C_WHITE)
        add_rect(sl, x, Inches(2.28), Inches(4.1), Inches(4.7),
                 RGBColor(0x2A, 0x1E, 0x12))
        for j, item in enumerate(p["items"]):
            add_textbox(sl, item,
                        x + Inches(0.18), Inches(2.4 + j * 0.86),
                        Inches(3.78), Inches(0.78), size=14, color=C_LIGHT)

    add_rect(sl, Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.4),
             RGBColor(0x3A, 0x2A, 0x08))
    add_textbox(sl,
                "90日後のゴール：物件契約完了・融資確定・SNS認知スタート　→　開業へ",
                Inches(0.55), Inches(7.07), Inches(12.3), Inches(0.36),
                size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    return sl


def s12_closing(prs):
    sl = blank_slide(prs)
    bg(sl)
    add_rect(sl, 0, 0, Inches(0.3), SLIDE_H, C_ACCENT)

    add_textbox(sl, "然", Inches(7.5), Inches(0.1), Inches(5.6), Inches(7.0),
                size=300, bold=True, color=RGBColor(0x2A, 0x1E, 0x12),
                align=PP_ALIGN.CENTER)

    add_textbox(sl, "私たちが目指す未来",
                Inches(0.9), Inches(1.1), Inches(6.8), Inches(0.65),
                size=18, bold=True, color=C_ACCENT)
    add_line(sl, Inches(0.9), Inches(1.82), Inches(5.7))

    add_textbox(sl,
                "糀漬け肉の定食と、和の一杯。",
                Inches(0.9), Inches(1.98), Inches(7.2), Inches(0.78),
                size=28, bold=True, color=C_WHITE)

    add_textbox(sl,
                "日本全国に「然」の灯りを。\n7年で100店舗。\n発酵と肉で、新しい定番をつくる。",
                Inches(0.9), Inches(2.9), Inches(6.8), Inches(1.7),
                size=22, color=C_LIGHT)

    add_textbox(sl,
                "肉酒場 然は、単なる飲食店ではありません。\n"
                "糀という日本の発酵文化と、現代人の「おいしくて、体にいい」\n"
                "という本質的なニーズを結びつけた、新しいブランドです。\n\n"
                "直営で型をつくり、FCで全国へ。\n"
                "その先には、日本の食文化に残る一皿があります。",
                Inches(0.9), Inches(4.75), Inches(6.8), Inches(2.1),
                size=15, color=C_GRAY, italic=True)

    slide_number(sl, 21)
    return sl


def s_group_strategy(prs, slide_num):
    """グループ事業戦略（3事業構成）"""
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "GROUP STRATEGY", "グループ事業戦略　――　3つの柱で成長する", slide_num)

    # 3事業カード
    businesses = [
        {
            "no": "01",
            "name": "肉酒場 然",
            "tag": "メイン事業　Year1〜",
            "icon": "🥩",
            "desc": "糀漬け肉×定食居酒屋の直営・FC展開。\n7年目100店舗・年商86億を目指すコア事業。",
            "y7": "年商 86億円",
            "kpi": "100店舗 / FC展開 / 上場検討",
        },
        {
            "no": "02",
            "name": "蕎麦 然",
            "tag": "立ち食い蕎麦FC　Year4〜",
            "icon": "🍜",
            "desc": "然の発酵だし×立ち食い蕎麦の\nFC業態。2人OP・客単価1,500円。\n月商400万円・高回転モデルで\nFC加盟者の収益性を高く設計。",
            "y7": "年商 12億円（約25店舗）",
            "kpi": "2人OP / 月商400万 / 客単価1,500円",
        },
        {
            "no": "03",
            "name": "SNS事業部",
            "tag": "事業部化　Year3〜",
            "icon": "📱",
            "desc": "3年目に事業部として独立。\nSNSコンサルティングと\nクリエイティブ制作を軸に\n飲食業界特化で展開。",
            "y7": "年商 2.5億円",
            "kpi": "月額顧問契約 / 案件単価UP",
        },
    ]

    for i, b in enumerate(businesses):
        x = Inches(0.55 + i * 4.25)
        # カード背景
        add_rect(sl, x, Inches(1.82), Inches(4.05), Inches(5.35),
                 RGBColor(0x2A, 0x1E, 0x12))
        # 番号・タグ
        add_textbox(sl, b["no"], x + Inches(0.15), Inches(1.95),
                    Inches(0.5), Inches(0.4), size=13, bold=True, color=C_ACCENT)
        add_rect(sl, x + Inches(0.6), Inches(1.98), Inches(3.2), Inches(0.32), C_ACCENT)
        add_textbox(sl, b["tag"], x + Inches(0.6), Inches(1.98),
                    Inches(3.2), Inches(0.32), size=11, bold=True,
                    color=C_BG, align=PP_ALIGN.CENTER)
        # 事業名
        add_textbox(sl, b["name"], x + Inches(0.15), Inches(2.38),
                    Inches(3.8), Inches(0.85), size=24, bold=True, color=C_WHITE)
        # 説明
        add_line(sl, x + Inches(0.15), Inches(3.28), Inches(3.7), color=C_GRAY,
                 height=Pt(0.5))
        add_textbox(sl, b["desc"], x + Inches(0.15), Inches(3.4),
                    Inches(3.75), Inches(1.3), size=14, color=C_LIGHT)
        # 7年目目標
        add_rect(sl, x, Inches(4.85), Inches(4.05), Inches(0.52),
                 RGBColor(0x3A, 0x2A, 0x08))
        add_textbox(sl, f"7年目目標：{b['y7']}", x + Inches(0.15), Inches(4.88),
                    Inches(3.8), Inches(0.42), size=15, bold=True, color=C_ACCENT)
        add_textbox(sl, b["kpi"], x + Inches(0.15), Inches(5.45),
                    Inches(3.8), Inches(0.6), size=12, color=C_GRAY)

    # 合計
    add_rect(sl, Inches(0.55), Inches(6.2), Inches(12.3), Inches(0.9),
             RGBColor(0x3A, 0x2A, 0x08))
    add_textbox(sl, "7年目グループ合計年商",
                Inches(0.8), Inches(6.28), Inches(4.0), Inches(0.38),
                size=16, color=C_GRAY)
    y7_total = REVENUE[6]
    add_textbox(sl, f"約{y7_total:.0f}億円",
                Inches(4.5), Inches(6.25), Inches(4.0), Inches(0.5),
                size=28, bold=True, color=C_ACCENT)
    add_textbox(sl, "（然86億 ＋ 蕎麦 然12億 ＋ SNS事業部4.5億）",
                Inches(8.3), Inches(6.33), Inches(4.3), Inches(0.38),
                size=13, color=C_GRAY)

    slide_number(sl, slide_num)
    return sl


def s_hq_cost(prs, slide_num):
    """本部経費・賞与引当スライド"""
    sl = blank_slide(prs)
    bg(sl)
    slide_header(sl, "HQ COST", "本部経費・賞与引当　――　スケールと共に整備する", slide_num)

    # 左：本部経費計画
    add_textbox(sl, "本部経費ロードマップ（Year3〜）",
                Inches(0.65), Inches(1.85), Inches(6.0), Inches(0.45),
                size=16, bold=True, color=C_ACCENT)

    hq_rows = [
        ("3年目",  "0.3億円",  "SV1名・経理外注・小オフィス"),
        ("4年目",  "0.8億円",  "SV2名・人事採用・本部スタッフ3名"),
        ("5年目",  "1.5億円",  "本部スタッフ6名・FC支援チーム"),
        ("6年目",  "2.5億円",  "SV5名・マーケ専任・IT担当"),
        ("7年目",  "4.0億円",  "本部組織完成形（100店舗対応）"),
        ("8〜10年目", "5〜7億円", "上場準備・海外展開準備コスト"),
    ]
    for i, (yr, cost, note) in enumerate(hq_rows):
        y = Inches(2.45 + i * 0.7)
        bg_c = RGBColor(0x3A, 0x2A, 0x08) if yr == "7年目" else RGBColor(0x2A, 0x1E, 0x12)
        add_rect(sl, Inches(0.55), y, Inches(6.3), Inches(0.65), bg_c)
        add_textbox(sl, yr,   Inches(0.7),  y + Inches(0.1), Inches(1.0), Inches(0.45),
                    size=14, bold=(yr == "7年目"), color=C_ACCENT)
        add_textbox(sl, cost, Inches(1.75), y + Inches(0.1), Inches(1.2), Inches(0.45),
                    size=16, bold=True, color=C_WHITE)
        add_textbox(sl, note, Inches(3.0),  y + Inches(0.12), Inches(3.7), Inches(0.42),
                    size=12, color=C_GRAY)

    # 右：商標料・システム料
    add_textbox(sl, "賞与引当（売上の1%相当）",
                Inches(7.3), Inches(1.85), Inches(5.7), Inches(0.45),
                size=16, bold=True, color=C_ACCENT)

    add_textbox(sl,
                "全店舗売上に対して1%を\nグループ本部へ計上。\n\n"
                "・ブランド維持費（商標登録・更新）\n"
                "・POSシステム・オペレーションツール\n"
                "・本部マーケティング共通費\n"
                "・FC加盟店サポート費用",
                Inches(7.3), Inches(2.45), Inches(5.7), Inches(2.2),
                size=15, color=C_LIGHT)

    add_rect(sl, Inches(7.1), Inches(4.8), Inches(5.9), Inches(0.55), C_ACCENT)
    add_textbox(sl, "7年目 賞与引当総額",
                Inches(7.2), Inches(4.83), Inches(3.5), Inches(0.44),
                size=14, bold=True, color=C_BG)
    add_textbox(sl, f"約{BONUS[6]:.1f}億円/年",
                Inches(10.6), Inches(4.83), Inches(2.2), Inches(0.44),
                size=18, bold=True, color=C_BG, align=PP_ALIGN.RIGHT)

    add_textbox(sl,
                "賞与引当は売上1%で毎月積立。社員・バイトへの賞与原資として管理。\n"
                "→ 100店舗時、本部だけで年商ベース約35億円規模のグループ収入",
                Inches(7.1), Inches(5.55), Inches(5.9), Inches(1.0),
                size=13, color=C_LIGHT, italic=True)

    # 利益率への影響
    add_rect(sl, Inches(0.55), Inches(6.75), Inches(12.3), Inches(0.55),
             RGBColor(0x2A, 0x1E, 0x12))
    add_textbox(sl,
                "本部経費・賞与引当控除後の実効営業利益率：Year3〜9% / Year7〜12% / Year10〜13%（15%から約2〜3%下方修正）",
                Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.44),
                size=13, bold=True, color=C_ACCENT)

    slide_number(sl, slide_num)
    return sl


def main():
    prs = new_prs()
    total = 23

    s01_title(prs)
    s02_vision(prs)
    s03_why_now(prs)
    s04_concept(prs)
    s05_products(prs)
    s06_competitor(prs)
    s07_profile(prs)
    s08_pl(prs)
    s09_finance(prs)
    s10_roadmap(prs)
    s_group_strategy(prs, slide_num=11)
    s11_strategy(prs)
    s13_open_strategy(prs)
    s14_fc_strategy(prs)
    s15_location(prs)
    s16_brand_sns(prs)
    s17_operations(prs)
    s18_data_kpi(prs)
    s_hq_cost(prs, slide_num=19)
    s_growth_chart(prs, slide_num=20)
    s19_risk(prs)
    s20_action_plan(prs)
    s12_closing(prs)

    os.makedirs("reports", exist_ok=True)
    out = f"reports/肉酒場然_事業計画書_{datetime.now().strftime('%Y%m%d')}.pptx"
    prs.save(out)
    print(f"保存完了: {out}")


if __name__ == "__main__":
    main()
