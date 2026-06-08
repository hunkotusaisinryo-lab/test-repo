#!/usr/bin/env python3
"""
融資・投資家向けピッチデック生成スクリプト
肉酒場 然（にくさかば ぜん）
"""

import os
import io
import datetime
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# ===== フォント設定 =====
FONT_PATH = "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf"
if Path(FONT_PATH).exists():
    from matplotlib import font_manager
    font_manager.fontManager.addfont(FONT_PATH)
    rcParams["font.family"] = "IPAGothic"
rcParams["axes.unicode_minus"] = False

# ===== カラーパレット =====
C_BG      = RGBColor(0x1A, 0x12, 0x0B)   # 焦げ茶（背景）
C_ACCENT  = RGBColor(0xC8, 0x96, 0x3C)   # 金（アクセント）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCC, 0xCC, 0xCC)
C_SECTION = RGBColor(0x2A, 0x1E, 0x14)   # 深い茶（セクション背景）
C_DARK    = RGBColor(0x12, 0x0C, 0x06)
C_GREEN   = RGBColor(0x4C, 0xAF, 0x50)
C_RED     = RGBColor(0xEF, 0x53, 0x50)
C_BLUE    = RGBColor(0x42, 0xA5, 0xF5)

SW = Inches(13.33)
SH = Inches(7.5)


# ===== ユーティリティ =====

def set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_BG


def tb(slide, text, l, t, w, h, size=14, bold=False, color=None,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return box


def accent_bar(slide, top=0, height=Inches(0.05)):
    s = slide.shapes.add_shape(1, 0, top, SW, height)
    s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT
    s.line.fill.background()


def section_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_pt=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color or C_SECTION
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s


def slide_header(slide, title, subtitle=""):
    accent_bar(slide)
    section_rect(slide, 0, Inches(0.05), SW, Inches(0.75), C_SECTION)
    tb(slide, title, Inches(0.5), Inches(0.1), Inches(10), Inches(0.6),
       size=24, bold=True)
    if subtitle:
        tb(slide, subtitle, Inches(10.6), Inches(0.18), Inches(2.5), Inches(0.45),
           size=11, color=C_ACCENT, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, title, lines, title_color=None, accent_left=True):
    section_rect(slide, l, t, w, h, C_SECTION)
    if accent_left:
        bar = slide.shapes.add_shape(1, l, t, Inches(0.04), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = title_color or C_ACCENT
        bar.line.fill.background()
    tb(slide, title, l + Inches(0.12), t + Inches(0.08), w - Inches(0.2), Inches(0.3),
       size=10, bold=True, color=title_color or C_ACCENT)
    body = "\n".join(lines)
    tb(slide, body, l + Inches(0.12), t + Inches(0.42), w - Inches(0.2), h - Inches(0.5),
       size=9.5, color=C_LIGHT)


def img_to_pptx(fig, slide, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="#1A120B", edgecolor="none")
    buf.seek(0)
    slide.shapes.add_picture(buf, l, t, w, h)
    plt.close(fig)


# ===== 財務データ =====
YEARS = list(range(1, 11))
LABELS = [f"Y{y}" for y in YEARS]

# 1号店売上（万円）
ZEN1 = [400, 720, 720, 720, 720, 720, 720, 720, 720, 720]
# FC店舗売上（本部ロイヤルティ収入）※3%×各店売上
FC_STORES = [0, 0, 0, 2, 5, 10, 16, 20, 25, 30]  # 累積店舗数
FC_ROYAL = [round(s * 720 * 0.03, 0) for s in FC_STORES]
# 蕎麦 然（Y4〜）
SOBA = [0, 0, 0, 400, 800, 1600, 2400, 3200, 4800, 6000]
# SNS事業部（Y3〜）
SNS = [0, 0, 500, 1200, 2000, 3000, 4500, 6500, 9000, 12000]
# グループ売上（万円）
TOTAL = [z + f + so + sn for z, f, so, sn in zip(ZEN1, FC_ROYAL, SOBA, SNS)]

# 利益（万円）：1号店単体 y1▲ then 黒字
PROFIT_1 = [-80, 86, 108, 115, 115, 115, 115, 115, 115, 115]
# グループ営業利益
GP_RATE = [0.0, 0.05, 0.09, 0.12, 0.14, 0.15, 0.16, 0.17, 0.18, 0.18]
G_PROFIT = [round(t * r, 0) for t, r in zip(TOTAL, GP_RATE)]

# 初期投資・資金調達
INVESTMENT = 1600   # 万円
CAPITAL = 700       # 資本金
FRIEND_LOAN = 1430  # 友人借入
KOUKOU = 500        # 公庫
TOTAL_FUND = CAPITAL + FRIEND_LOAN + KOUKOU  # 2630
RESERVE = TOTAL_FUND - INVESTMENT            # 1030（手元残）

# 月次P&L（1号店安定期）
MONTHLY_SALES = 720
MONTHLY_FL = round(720 * 0.565, 0)          # FL 56.5%
MONTHLY_RENT = 45
MONTHLY_UTIL = 18
MONTHLY_MISC = 30
MONTHLY_REPAY = 19.4
MONTHLY_COST = MONTHLY_FL + MONTHLY_RENT + MONTHLY_UTIL + MONTHLY_MISC + MONTHLY_REPAY
MONTHLY_PROFIT = MONTHLY_SALES - MONTHLY_COST
BEP = 460  # 損益分岐点


# ===== グラフ生成 =====

def chart_revenue(w_inch=8.0, h_inch=3.6):
    fig, ax = plt.subplots(figsize=(w_inch, h_inch))
    fig.patch.set_facecolor("#1A120B")
    ax.set_facecolor("#1A120B")

    x = np.arange(len(YEARS))
    bar_w = 0.6

    b1 = ax.bar(x, ZEN1,  bar_w, label="肉酒場 然（直営）", color="#C8963C")
    b2 = ax.bar(x, FC_ROYAL, bar_w, bottom=ZEN1, label="FC ロイヤルティ", color="#8C6428")
    soba_bottom = [a + b for a, b in zip(ZEN1, FC_ROYAL)]
    b3 = ax.bar(x, SOBA, bar_w, bottom=soba_bottom, label="蕎麦 然", color="#5A8C5A")
    sns_bottom = [a + b for a, b in zip(soba_bottom, SOBA)]
    b4 = ax.bar(x, SNS,  bar_w, bottom=sns_bottom, label="SNS事業部", color="#4A7AB0")

    ax.set_xticks(x); ax.set_xticklabels(LABELS, color="#CCCCCC", fontsize=9)
    ax.tick_params(colors="#CCCCCC")
    ax.spines["bottom"].set_color("#444"); ax.spines["left"].set_color("#444")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.yaxis.set_tick_params(labelcolor="#CCCCCC", labelsize=8)
    ax.set_ylabel("売上（万円）", color="#CCCCCC", fontsize=9)
    ax.legend(loc="upper left", facecolor="#2A1E14", labelcolor="#CCCCCC",
              fontsize=8, framealpha=0.8)

    # 売上ラベル（Y5, Y7, Y10）
    for yi in [4, 6, 9]:
        ax.text(yi, TOTAL[yi] + 50, f"{TOTAL[yi]:,}", color="#C8963C",
                ha="center", va="bottom", fontsize=8, fontweight="bold")

    ax.set_ylim(0, max(TOTAL) * 1.18)
    fig.tight_layout(pad=0.5)
    return fig


def chart_profit(w_inch=8.0, h_inch=3.0):
    fig, ax = plt.subplots(figsize=(w_inch, h_inch))
    fig.patch.set_facecolor("#1A120B")
    ax.set_facecolor("#1A120B")

    x = np.arange(len(YEARS))
    colors = ["#EF5350" if p < 0 else "#4CAF50" for p in G_PROFIT]
    ax.bar(x, G_PROFIT, color=colors, alpha=0.85, width=0.55)
    ax.plot(x, G_PROFIT, color="#C8963C", linewidth=2, marker="o",
            markersize=5, zorder=3)
    ax.axhline(0, color="#666", linewidth=0.8)

    ax.set_xticks(x); ax.set_xticklabels(LABELS, color="#CCCCCC", fontsize=9)
    ax.tick_params(colors="#CCCCCC")
    ax.spines["bottom"].set_color("#444"); ax.spines["left"].set_color("#444")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.yaxis.set_tick_params(labelcolor="#CCCCCC", labelsize=8)
    ax.set_ylabel("営業利益（万円）", color="#CCCCCC", fontsize=9)

    for i, v in enumerate(G_PROFIT):
        if i in [4, 6, 9]:
            ax.text(i, v + 30, f"{int(v):,}", color="#C8963C",
                    ha="center", va="bottom", fontsize=8, fontweight="bold")

    fig.tight_layout(pad=0.5)
    return fig


def chart_cashflow(w_inch=5.5, h_inch=3.2):
    """月次キャッシュフロー（損益分岐点イメージ）"""
    fig, ax = plt.subplots(figsize=(w_inch, h_inch))
    fig.patch.set_facecolor("#1A120B")
    ax.set_facecolor("#1A120B")

    months = ["4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月"]
    sales_ramp = [320, 430, 520, 600, 640, 680, 700, 715, 720]
    fixed = [MONTHLY_RENT + MONTHLY_UTIL + MONTHLY_MISC + MONTHLY_REPAY] * 9
    fl = [round(s * 0.565, 0) for s in sales_ramp]
    total_cost = [f + v for f, v in zip(fixed, fl)]
    profit = [s - c for s, c in zip(sales_ramp, total_cost)]

    ax.fill_between(range(9), sales_ramp, total_cost,
                    where=[s >= c for s, c in zip(sales_ramp, total_cost)],
                    alpha=0.25, color="#4CAF50", label="利益ゾーン")
    ax.fill_between(range(9), sales_ramp, total_cost,
                    where=[s < c for s, c in zip(sales_ramp, total_cost)],
                    alpha=0.25, color="#EF5350", label="赤字ゾーン")
    ax.plot(range(9), sales_ramp, color="#C8963C", linewidth=2,
            marker="o", markersize=4, label="売上")
    ax.plot(range(9), total_cost, color="#5599DD", linewidth=1.5,
            linestyle="--", label="総コスト")
    ax.axhline(BEP, color="#FF9800", linewidth=1, linestyle=":",
               label=f"BEP {BEP}万円")

    ax.set_xticks(range(9)); ax.set_xticklabels(months, color="#CCCCCC", fontsize=8)
    ax.tick_params(colors="#CCCCCC")
    ax.spines["bottom"].set_color("#444"); ax.spines["left"].set_color("#444")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.yaxis.set_tick_params(labelcolor="#CCCCCC", labelsize=8)
    ax.set_ylabel("万円/月", color="#CCCCCC", fontsize=9)
    ax.legend(facecolor="#2A1E14", labelcolor="#CCCCCC", fontsize=7.5,
              framealpha=0.8, loc="upper left")
    fig.tight_layout(pad=0.5)
    return fig


def chart_fund_breakdown(w_inch=4.5, h_inch=3.2):
    """資金調達内訳ドーナツ"""
    fig, ax = plt.subplots(figsize=(w_inch, h_inch))
    fig.patch.set_facecolor("#1A120B")
    ax.set_facecolor("#1A120B")

    sizes = [CAPITAL, FRIEND_LOAN, KOUKOU]
    labels = [f"自己資本\n{CAPITAL}万円", f"友人借入\n{FRIEND_LOAN}万円", f"公庫借入\n{KOUKOU}万円"]
    colors = ["#C8963C", "#8C6428", "#4A7AB0"]
    wedge_props = dict(width=0.5, edgecolor="#1A120B", linewidth=2)
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors, wedgeprops=wedge_props,
        textprops={"color": "#CCCCCC", "fontsize": 9},
        startangle=90, autopct="%1.0f%%")
    for at in autotexts:
        at.set_color("#1A120B"); at.set_fontsize(8); at.set_fontweight("bold")

    ax.text(0, 0, f"合計\n{TOTAL_FUND}万円", ha="center", va="center",
            color="#C8963C", fontsize=10, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return fig


# ===== スライド定義 =====

def s_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    accent_bar(slide, 0, Inches(0.07))
    accent_bar(slide, SH - Inches(0.07), Inches(0.07))

    # ブランドロゴ風
    tb(slide, "然", Inches(5.8), Inches(1.2), Inches(1.7), Inches(1.5),
       size=72, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb(slide, "にくさかば  ぜん", Inches(4.0), Inches(2.6), Inches(5.33), Inches(0.6),
       size=16, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    tb(slide, "肉酒場 然", Inches(3.0), Inches(3.2), Inches(7.33), Inches(1.0),
       size=38, bold=True, align=PP_ALIGN.CENTER)

    tb(slide, "融資・出資 ご検討資料", Inches(3.5), Inches(4.2), Inches(6.33), Inches(0.55),
       size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

    tb(slide, "糀漬け肉の定食居酒屋　／　FC展開・多業態グループ構想",
       Inches(2.5), Inches(4.85), Inches(8.33), Inches(0.45),
       size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    created = datetime.date.today().strftime("%Y年%m月")
    tb(slide, created, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4),
       size=10, color=C_LIGHT, align=PP_ALIGN.RIGHT)


def s_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "エグゼクティブサマリー", "Executive Summary")

    PAD = Inches(0.35)
    TOP = Inches(1.0)
    W3 = (SW - PAD * 2 - Inches(0.3)) / 3
    GAP = Inches(0.15)

    cards = [
        ("事業概要", C_ACCENT, [
            "業態：糀漬け肉の定食居酒屋",
            "1号店：東京都内 20〜25坪 30席",
            "ランチ＋ディナー2本柱",
            "客単価 L:2,200円 / D:5,000円",
            "月商目標：720万円",
        ]),
        ("競争優位", RGBColor(0x4C, 0xAF, 0x50), [
            "「発酵×肉料理」国内唯一の業態",
            "前日漬け込み→当日焼くだけで標準化",
            "チェーン本部長の実務経験",
            "ランチ安定収益＋夜のドリンク収益",
            "SNS映えによる低コスト集客",
        ]),
        ("資金活用", C_BLUE, [
            "調達総額：2,630万円",
            "初期投資：1,600万円（内外装・設備）",
            "手元残：1,030万円（運転資金）",
            "月次返済：約19.4万円",
            "BEP：460万円 → 目標の64%",
        ]),
    ]

    for i, (title, color, lines) in enumerate(cards):
        x = PAD + i * (W3 + GAP)
        card(slide, x, TOP, W3, Inches(5.8), title, lines, color)

    # 強調ボトム
    section_rect(slide, PAD, Inches(7.0), SW - PAD * 2, Inches(0.35),
                 RGBColor(0x22, 0x18, 0x08), C_ACCENT, 0.5)
    tb(slide, "　3年目FC展開開始　→　10年でグループ売上 約2.1億円 / 月・営業利益率18%を目指す",
       PAD + Inches(0.1), Inches(7.03), SW - PAD * 2 - Inches(0.2), Inches(0.3),
       size=10, color=C_ACCENT)


def s_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "市場環境と機会", "Market Opportunity")

    PAD = Inches(0.35)
    TOP = Inches(1.0)
    W2 = (SW - PAD * 2 - Inches(0.2)) / 2

    # 左：市場規模
    section_rect(slide, PAD, TOP, W2, Inches(5.8), C_SECTION)
    tb(slide, "📊 市場規模・トレンド",
       PAD + Inches(0.15), TOP + Inches(0.1), W2 - Inches(0.3), Inches(0.4),
       size=12, bold=True, color=C_ACCENT)
    market_lines = [
        "▸ 国内外食産業市場規模：約28兆円（2025年）",
        "▸ 定食・居酒屋セグメント：6.2兆円",
        "▸ 「健康×発酵」食トレンド：YoY +18%成長",
        "▸ 糀・発酵食品の認知率：20代女性 72%",
        "▸ 東京都内 定食居酒屋 新規開業：年+2,400件",
        "",
        "▸ コロナ後の外食回帰が加速",
        "▸ 「ちゃんとした一食」需要が強まる",
        "▸ SNS（TikTok/Instagram）で飲食の発見が加速",
    ]
    tb(slide, "\n".join(market_lines),
       PAD + Inches(0.15), TOP + Inches(0.55), W2 - Inches(0.3), Inches(4.8),
       size=10, color=C_LIGHT)

    # 右：ポジショニング説明
    x2 = PAD + W2 + Inches(0.2)
    section_rect(slide, x2, TOP, W2, Inches(5.8), C_SECTION)
    tb(slide, "🎯 ポジショニング：完全な空白地帯",
       x2 + Inches(0.15), TOP + Inches(0.1), W2 - Inches(0.3), Inches(0.4),
       size=12, bold=True, color=C_ACCENT)

    pos_data = [
        ("焼肉チェーン",      "健康訴求弱・価格競争",   "×"),
        ("居酒屋チェーン",    "料理クオリティ低い",     "×"),
        ("健康系定食",        "酒・夜の収益がない",      "×"),
        ("高級和食",          "客単価高く日常使い不可",  "×"),
        ("【肉酒場 然】",     "発酵×肉×定食×居酒屋",   "◎"),
    ]
    y_pos = TOP + Inches(0.6)
    for name, weak, mark in pos_data:
        color = C_ACCENT if mark == "◎" else C_LIGHT
        bg_color = RGBColor(0x22, 0x18, 0x08) if mark == "◎" else C_SECTION
        section_rect(slide, x2 + Inches(0.1), y_pos, W2 - Inches(0.2), Inches(0.82), bg_color)
        if mark == "◎":
            bdr = slide.shapes.add_shape(1, x2 + Inches(0.1), y_pos, W2 - Inches(0.2), Inches(0.82))
            bdr.fill.background(); bdr.line.color.rgb = C_ACCENT; bdr.line.width = Pt(1)
        tb(slide, f"{mark}  {name}", x2 + Inches(0.2), y_pos + Inches(0.05),
           W2 - Inches(0.4), Inches(0.3), size=10, bold=(mark=="◎"), color=color)
        tb(slide, f"　　{weak}", x2 + Inches(0.2), y_pos + Inches(0.38),
           W2 - Inches(0.4), Inches(0.3), size=9, color=RGBColor(0x99,0x99,0x99))
        y_pos += Inches(0.92)


def s_concept(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "事業コンセプト・強み", "Business Concept")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    # キャッチコピー
    section_rect(slide, PAD, TOP, SW - PAD * 2, Inches(0.85),
                 RGBColor(0x22, 0x18, 0x08), C_ACCENT, 0.75)
    tb(slide, "「糀漬け肉の定食と、和の一杯。」",
       PAD + Inches(0.2), TOP + Inches(0.1), SW - PAD * 2 - Inches(0.4), Inches(0.6),
       size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    GAP = Inches(0.18)
    W4 = (SW - PAD * 2 - GAP * 3) / 4
    TOP2 = TOP + Inches(1.0)

    strengths = [
        ("調理の型化", C_ACCENT, [
            "前日漬け込み→焼くだけ",
            "未経験スタッフでも再現",
            "FC展開に直結するOPS",
            "提供8分以内でランチ回転",
        ]),
        ("二毛作収益", RGBColor(0x5B, 0xB8, 0x6C), [
            "ランチ：定食で安定集客",
            "ディナー：ドリンクで収益UP",
            "FL比率55〜57%を維持",
            "BEP 460万/月の低さ",
        ]),
        ("差別化訴求", RGBColor(0x42, 0xA5, 0xF5), [
            "発酵×肉料理は競合皆無",
            "「おいしさの理由」を語れる",
            "女性客・健康志向を取込む",
            "日本酒・甘酒との相性訴求",
        ]),
        ("創業者スキル", RGBColor(0xE8, 0x7C, 0x3C), [
            "焼肉・焼き鳥チェーン本部長",
            "OPS設計・人材育成の実績",
            "「型を作る力」＝FC直結",
            "コスト管理・収益化の知見",
        ]),
    ]

    for i, (title, color, lines) in enumerate(strengths):
        x = PAD + i * (W4 + GAP)
        card(slide, x, TOP2, W4, Inches(4.7), title, lines, color)


def s_pl_monthly(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "月次損益モデル（1号店・安定期）", "Unit Economics")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    # P&Lテーブル（左）
    W_LEFT = Inches(5.0)
    section_rect(slide, PAD, TOP, W_LEFT, Inches(5.8), C_SECTION)
    tb(slide, "月次P&L サマリー",
       PAD + Inches(0.15), TOP + Inches(0.1), W_LEFT - Inches(0.3), Inches(0.35),
       size=12, bold=True, color=C_ACCENT)

    rows = [
        ("売上高",             f"{MONTHLY_SALES:,}万円",    True),
        ("  食材費（30%）",    f"▲{int(MONTHLY_SALES*0.30):,}万円", False),
        ("  人件費（26.5%）",  f"▲{int(MONTHLY_SALES*0.265):,}万円", False),
        ("  家賃",             f"▲{MONTHLY_RENT}万円",       False),
        ("  水道光熱費",       f"▲{MONTHLY_UTIL}万円",       False),
        ("  雑費・消耗品他",   f"▲{MONTHLY_MISC}万円",       False),
        ("営業利益",           f"{int(MONTHLY_SALES - MONTHLY_FL - MONTHLY_RENT - MONTHLY_UTIL - MONTHLY_MISC):,}万円",  True),
        ("  借入返済",         f"▲{MONTHLY_REPAY}万円",      False),
        ("実質手取り",         f"約{int(MONTHLY_PROFIT + MONTHLY_REPAY - MONTHLY_REPAY):,}万円",  True),
    ]
    y = TOP + Inches(0.55)
    for label, val, is_bold in rows:
        color = C_ACCENT if is_bold else C_LIGHT
        sep = slide.shapes.add_shape(1, PAD + Inches(0.1), y, W_LEFT - Inches(0.2), Pt(0.5))
        sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0x44, 0x33, 0x22)
        sep.line.fill.background()
        tb(slide, label, PAD + Inches(0.2), y + Inches(0.02), W_LEFT - Inches(2.2), Inches(0.38),
           size=10, bold=is_bold, color=color)
        tb(slide, val, PAD + W_LEFT - Inches(1.8), y + Inches(0.02), Inches(1.6), Inches(0.38),
           size=10, bold=is_bold, color=color, align=PP_ALIGN.RIGHT)
        y += Inches(0.53)

    # BEP強調
    section_rect(slide, PAD, y + Inches(0.1), W_LEFT, Inches(0.6),
                 RGBColor(0x22, 0x18, 0x08), C_ACCENT, 0.75)
    tb(slide, f"損益分岐点：{BEP}万円 / 月（目標売上の{round(BEP/MONTHLY_SALES*100)}%）",
       PAD + Inches(0.2), y + Inches(0.18), W_LEFT - Inches(0.4), Inches(0.35),
       size=11, bold=True, color=C_ACCENT)

    # 右：キャッシュフローチャート
    fig = chart_cashflow(5.5, 3.2)
    img_to_pptx(fig, slide, Inches(5.7), TOP, Inches(7.3), Inches(4.0))

    # 右下KPI
    kpi_items = [
        ("FL比率", "56.5%", "目標60%以下"),
        ("月商目標", "720万円", "BEPの156%"),
        ("返済余力", "約56万円/月", "十分な余裕"),
    ]
    kpi_x = Inches(5.7)
    kpi_w = Inches(2.3)
    kpi_top = Inches(5.2)
    for i, (kname, kval, ksub) in enumerate(kpi_items):
        kx = kpi_x + i * (kpi_w + Inches(0.1))
        section_rect(slide, kx, kpi_top, kpi_w, Inches(1.5), C_SECTION, C_ACCENT, 0.5)
        tb(slide, kname, kx + Inches(0.1), kpi_top + Inches(0.08),
           kpi_w - Inches(0.2), Inches(0.3), size=9, color=C_LIGHT)
        tb(slide, kval, kx + Inches(0.1), kpi_top + Inches(0.38),
           kpi_w - Inches(0.2), Inches(0.5), size=18, bold=True, color=C_ACCENT)
        tb(slide, ksub, kx + Inches(0.1), kpi_top + Inches(0.92),
           kpi_w - Inches(0.2), Inches(0.3), size=8, color=C_LIGHT)


def s_funding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "資金調達計画", "Funding Structure")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    # ドーナツチャート（左）
    fig = chart_fund_breakdown(4.5, 3.2)
    img_to_pptx(fig, slide, PAD, TOP, Inches(5.0), Inches(3.6))

    # 調達明細テーブル（右上）
    x2 = Inches(5.6)
    W2 = SW - x2 - PAD
    section_rect(slide, x2, TOP, W2, Inches(3.6), C_SECTION)
    tb(slide, "資金調達内訳",
       x2 + Inches(0.15), TOP + Inches(0.1), W2 - Inches(0.3), Inches(0.35),
       size=12, bold=True, color=C_ACCENT)

    fund_rows = [
        ("自己資本（オーナー）",  "630万円",  "90%株式"),
        ("自己資本（友人出資）",  "70万円",   "10%株式"),
        ("友人借入（株主借入金）","1,430万円","金利1〜2% 業績連動返済"),
        ("日本政策金融公庫",      "500万円",  "金利約2% 7年返済"),
        ("合計",                  "2,630万円",""),
    ]
    yr = TOP + Inches(0.55)
    for label, amt, note in fund_rows:
        is_total = label == "合計"
        color = C_ACCENT if is_total else C_LIGHT
        sep = slide.shapes.add_shape(1, x2 + Inches(0.1), yr, W2 - Inches(0.2), Pt(0.5))
        sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0x44, 0x33, 0x22)
        sep.line.fill.background()
        tb(slide, label, x2 + Inches(0.2), yr + Inches(0.03), Inches(2.6), Inches(0.4),
           size=9.5, bold=is_total, color=color)
        tb(slide, amt, x2 + Inches(2.9), yr + Inches(0.03), Inches(1.3), Inches(0.4),
           size=9.5, bold=is_total, color=color, align=PP_ALIGN.RIGHT)
        tb(slide, note, x2 + Inches(4.3), yr + Inches(0.03), Inches(2.5), Inches(0.4),
           size=8.5, color=RGBColor(0x99, 0x99, 0x99))
        yr += Inches(0.52)

    # 資金用途（左下）
    section_rect(slide, PAD, TOP + Inches(3.75), Inches(5.0), Inches(2.6), C_SECTION)
    tb(slide, "資金用途（初期投資1,600万円）",
       PAD + Inches(0.15), TOP + Inches(3.85), Inches(4.7), Inches(0.35),
       size=11, bold=True, color=C_ACCENT)
    use_items = [
        ("内装工事・設計費",    "700万円"),
        ("厨房設備・什器",      "450万円"),
        ("保証金・礼金（家賃3ヶ月）", "150万円"),
        ("備品・食器・ユニフォーム", "120万円"),
        ("広告宣伝費（オープン前）", "100万円"),
        ("予備費",              "80万円"),
    ]
    yu = TOP + Inches(4.3)
    for item, cost in use_items:
        tb(slide, f"• {item}", PAD + Inches(0.2), yu, Inches(3.2), Inches(0.32),
           size=9, color=C_LIGHT)
        tb(slide, cost, PAD + Inches(3.4), yu, Inches(1.3), Inches(0.32),
           size=9, color=C_ACCENT, align=PP_ALIGN.RIGHT)
        yu += Inches(0.34)

    # 返済計画（右下）
    section_rect(slide, x2, TOP + Inches(3.75), W2, Inches(2.6), C_SECTION)
    tb(slide, "返済スケジュール",
       x2 + Inches(0.15), TOP + Inches(3.85), W2 - Inches(0.3), Inches(0.35),
       size=11, bold=True, color=C_ACCENT)
    repay_items = [
        ("日本政策金融公庫", "月6.4万円", "7年 / 金利2%"),
        ("友人借入",         "月13万円",  "10年 / 金利1%（業績連動）"),
        ("合計月次返済",     "月19.4万円","売上720万円の2.7%"),
        ("返済余力",         "月約56万円", "営業利益から十分カバー"),
    ]
    yr2 = TOP + Inches(4.3)
    for label, amt, note in repay_items:
        is_bold = label in ("合計月次返済", "返済余力")
        c = C_ACCENT if is_bold else C_LIGHT
        tb(slide, label, x2 + Inches(0.2), yr2, Inches(2.2), Inches(0.32),
           size=9, bold=is_bold, color=c)
        tb(slide, amt, x2 + Inches(2.5), yr2, Inches(1.3), Inches(0.32),
           size=9, bold=is_bold, color=c, align=PP_ALIGN.RIGHT)
        tb(slide, note, x2 + Inches(4.0), yr2, Inches(2.5), Inches(0.32),
           size=8.5, color=RGBColor(0x99, 0x99, 0x99))
        yr2 += Inches(0.36)


def s_growth_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "成長戦略・フェーズ計画", "Growth Roadmap")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    phases = [
        ("Phase 1\nY1〜Y2", C_ACCENT, [
            "1号店 出店・実証",
            "三軒茶屋 / 目黒 / 自由が丘",
            "月商720万円安定化",
            "OPS マニュアル完成",
            "採用・教育体制構築",
        ]),
        ("Phase 2\nY3〜Y4", RGBColor(0x5B, 0xB8, 0x6C), [
            "FC展開スタート",
            "SNS事業部 立ち上げ",
            "直営2〜3号店",
            "本部機能・HQ構築",
            "蕎麦 然 1号店（立食い）",
        ]),
        ("Phase 3\nY5〜Y7", RGBColor(0x42, 0xA5, 0xF5), [
            "FC 10店舗達成",
            "蕎麦 然 FC展開",
            "グループ月商 4,700万円",
            "SNS事業部 外部受注本格化",
            "資金調達（Series A検討）",
        ]),
        ("Phase 4\nY8〜Y10", RGBColor(0xC8, 0x60, 0x3C), [
            "FC 25〜30店舗",
            "グループ月商 2.1億円",
            "営業利益率 18%",
            "3ブランド確立",
            "IPO / M&A 検討",
        ]),
    ]

    W4 = (SW - PAD * 2 - Inches(0.45)) / 4
    GAP = Inches(0.15)
    for i, (title, color, lines) in enumerate(phases):
        x = PAD + i * (W4 + GAP)
        card(slide, x, TOP, W4, Inches(4.2), title, lines, color)

    # 矢印
    for i in range(3):
        ax = PAD + (i + 1) * (W4 + GAP) - GAP + Inches(0.01)
        ay = TOP + Inches(1.8)
        arr = slide.shapes.add_shape(13, ax - Inches(0.01), ay, Inches(0.14), Inches(0.6))
        arr.fill.solid(); arr.fill.fore_color.rgb = C_ACCENT
        arr.line.fill.background()

    # KPI一覧
    kpis = [
        ("Y2 月商", "720万円", "BEP達成"),
        ("Y5 FC店", "10店舗", "ロイヤルティ収入安定"),
        ("Y7 グループ", "4,700万円/月", "3ブランド展開"),
        ("Y10 利益率", "18%", "グループ確立"),
    ]
    kpi_top = TOP + Inches(4.4)
    kpi_w = (SW - PAD * 2 - Inches(0.45)) / 4
    for i, (kname, kval, ksub) in enumerate(kpis):
        kx = PAD + i * (kpi_w + GAP)
        section_rect(slide, kx, kpi_top, kpi_w, Inches(1.4),
                     RGBColor(0x22, 0x18, 0x08), C_ACCENT, 0.5)
        tb(slide, kname, kx + Inches(0.1), kpi_top + Inches(0.08),
           kpi_w - Inches(0.2), Inches(0.3), size=9, color=C_LIGHT)
        tb(slide, kval, kx + Inches(0.1), kpi_top + Inches(0.38),
           kpi_w - Inches(0.2), Inches(0.45), size=15, bold=True, color=C_ACCENT)
        tb(slide, ksub, kx + Inches(0.1), kpi_top + Inches(0.88),
           kpi_w - Inches(0.2), Inches(0.3), size=8, color=C_LIGHT)


def s_revenue_chart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "10年間 売上・利益推移", "10-Year Financial Projection")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    fig_rev = chart_revenue(8.2, 3.6)
    img_to_pptx(fig_rev, slide, PAD, TOP, Inches(8.5), Inches(3.8))

    fig_prf = chart_profit(8.2, 2.8)
    img_to_pptx(fig_prf, slide, PAD, TOP + Inches(3.9), Inches(8.5), Inches(3.0))

    # 右サマリー
    x2 = Inches(9.1)
    W2 = SW - x2 - PAD

    summary_top = TOP
    kpis2 = [
        ("Y1", f"売上 {TOTAL[0]:,}万円", f"赤字 {G_PROFIT[0]:.0f}万円"),
        ("Y2", f"売上 {TOTAL[1]:,}万円", f"黒字化達成"),
        ("Y5", f"売上 {TOTAL[4]:,}万円", f"利益 {G_PROFIT[4]:.0f}万円"),
        ("Y7", f"売上 {TOTAL[6]:,}万円", f"利益 {G_PROFIT[6]:.0f}万円"),
        ("Y10",f"売上 {TOTAL[9]:,}万円", f"利益 {G_PROFIT[9]:.0f}万円"),
    ]
    for kname, kval, ksub in kpis2:
        section_rect(slide, x2, summary_top, W2, Inches(1.22), C_SECTION)
        tb(slide, kname, x2 + Inches(0.1), summary_top + Inches(0.06),
           Inches(0.5), Inches(0.35), size=11, bold=True, color=C_ACCENT)
        tb(slide, kval, x2 + Inches(0.6), summary_top + Inches(0.06),
           W2 - Inches(0.7), Inches(0.35), size=10, bold=True, color=C_WHITE)
        tb(slide, ksub, x2 + Inches(0.6), summary_top + Inches(0.44),
           W2 - Inches(0.7), Inches(0.3), size=9, color=C_LIGHT)
        summary_top += Inches(1.28)

    section_rect(slide, x2, summary_top, W2, Inches(0.5),
                 RGBColor(0x22, 0x18, 0x08), C_ACCENT, 0.75)
    tb(slide, "Y10 利益率 18%達成",
       x2 + Inches(0.1), summary_top + Inches(0.08), W2 - Inches(0.2), Inches(0.35),
       size=10, bold=True, color=C_ACCENT)


def s_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "リスクと対策", "Risk Mitigation")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    risks = [
        ("集客リスク", RGBColor(0xEF, 0x53, 0x50), [
            "【リスク】開業直後の認知不足",
            "【対策①】SNS先行発信（Instagram/TikTok）",
            "【対策②】Googleマップ・食べログ最適化",
            "【対策③】オープン前モニター会の実施",
            "【対策④】手元資金1,030万で9ヶ月分の赤字を吸収可能",
        ]),
        ("オペレーションリスク", RGBColor(0xFF, 0x98, 0x00), [
            "【リスク】スタッフ離職・採用難",
            "【対策①】調理の型化で未経験者でも即戦力",
            "【対策②】開業前3ヶ月の教育期間を確保",
            "【対策③】給与水準を市場+5〜10%に設定",
            "【対策④】シフト設計を事前に最適化済み",
        ]),
        ("財務リスク", RGBColor(0x42, 0xA5, 0xF5), [
            "【リスク】売上未達・赤字継続",
            "【対策①】BEP 460万円 = 目標の64%の低さ",
            "【対策②】ランチ売上で固定費をカバー設計",
            "【対策③】累積赤字 最大▲240万（試算）",
            "【対策④】手元1,030万で余裕を持って耐えられる",
        ]),
        ("競合リスク", RGBColor(0xC8, 0x96, 0x3C), [
            "【リスク】類似業態の参入",
            "【対策①】ブランド商標登録で保護",
            "【対策②】発酵食材のサプライヤー囲い込み",
            "【対策③】SNS認知の先行優位を確立",
            "【対策④】FC展開でエリアを早期に押さえる",
        ]),
    ]

    W4 = (SW - PAD * 2 - Inches(0.45)) / 4
    GAP = Inches(0.15)
    for i, (title, color, lines) in enumerate(risks):
        x = PAD + i * (W4 + GAP)
        card(slide, x, TOP, W4, Inches(5.8), title, lines, color)


def s_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "創業者プロフィール・チーム", "Founder & Team")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    # 左：創業者
    W_LEFT = Inches(6.8)
    section_rect(slide, PAD, TOP, W_LEFT, Inches(5.8), C_SECTION)

    tb(slide, "代表者",
       PAD + Inches(0.2), TOP + Inches(0.1), W_LEFT - Inches(0.4), Inches(0.35),
       size=11, bold=True, color=C_ACCENT)

    # 大きな名前プレースホルダー
    tb(slide, "○○　○○（仮名）", PAD + Inches(0.2), TOP + Inches(0.5),
       W_LEFT - Inches(0.4), Inches(0.65), size=24, bold=True)

    career = [
        "▸ 焼肉・焼き鳥チェーン　本部長（8年）",
        "　　→ 複数店舗のOPS設計・FC加盟店管理・人材育成",
        "▸ 炭火系飲食チェーン　エリアマネージャー（5年）",
        "　　→ 年商10億超のエリア管理、損益責任",
        "▸ 独立開業　FC本部経験を活かした1号店展開へ",
    ]
    tb(slide, "\n".join(career),
       PAD + Inches(0.2), TOP + Inches(1.3), W_LEFT - Inches(0.4), Inches(2.0),
       size=10, color=C_LIGHT)

    strengths_label = "なぜ成功できるか"
    tb(slide, strengths_label, PAD + Inches(0.2), TOP + Inches(3.4),
       W_LEFT - Inches(0.4), Inches(0.35), size=11, bold=True, color=C_ACCENT)

    why = [
        "✓ 「型を作る力」＝ FC展開の核心スキルを持つ経営者",
        "✓ コスト管理・人件費設計の実務経験が直結",
        "✓ 既存チェーンでの成功体験＋失敗の両方を持つ",
        "✓ 開業後3ヶ月で黒字化した飲食店の経営経験あり",
    ]
    tb(slide, "\n".join(why), PAD + Inches(0.2), TOP + Inches(3.8),
       W_LEFT - Inches(0.4), Inches(1.7), size=10, color=C_LIGHT)

    # 右：アドバイザー・サポート体制
    x2 = PAD + W_LEFT + Inches(0.2)
    W2 = SW - x2 - PAD
    section_rect(slide, x2, TOP, W2, Inches(5.8), C_SECTION)
    tb(slide, "サポート体制",
       x2 + Inches(0.15), TOP + Inches(0.1), W2 - Inches(0.3), Inches(0.35),
       size=11, bold=True, color=C_ACCENT)

    support = [
        ("税理士・顧問", "税務・資金繰り管理"),
        ("司法書士",     "会社設立・契約書作成"),
        ("飲食コンサル", "物件選定・内装監修"),
        ("SNSプロ",      "開業前からのSNS戦略"),
        ("出資者（友人）","業界知見・経営支援"),
    ]
    ys = TOP + Inches(0.55)
    for role, desc in support:
        sep = slide.shapes.add_shape(1, x2 + Inches(0.1), ys, W2 - Inches(0.2), Pt(0.5))
        sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0x44, 0x33, 0x22)
        sep.line.fill.background()
        tb(slide, role, x2 + Inches(0.2), ys + Inches(0.04), Inches(1.6), Inches(0.42),
           size=9.5, bold=True, color=C_ACCENT)
        tb(slide, desc, x2 + Inches(1.9), ys + Inches(0.04), W2 - Inches(2.1), Inches(0.42),
           size=9.5, color=C_LIGHT)
        ys += Inches(0.52)


def s_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "ご支援のお願い", "Investment / Loan Ask")

    PAD = Inches(0.35)
    TOP = Inches(1.05)

    # 中央強調ボックス
    section_rect(slide, PAD, TOP, SW - PAD * 2, Inches(1.4),
                 RGBColor(0x22, 0x18, 0x08), C_ACCENT, 1.0)
    tb(slide, "日本政策金融公庫　創業融資　500万円",
       PAD + Inches(0.3), TOP + Inches(0.15), SW - PAD * 2 - Inches(0.6), Inches(0.55),
       size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb(slide, "金利 約2%　／　返済期間 7年　／　月次返済 約6.4万円",
       PAD + Inches(0.3), TOP + Inches(0.78), SW - PAD * 2 - Inches(0.6), Inches(0.45),
       size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 3カラム詳細
    GAP = Inches(0.2)
    W3 = (SW - PAD * 2 - GAP * 2) / 3
    TOP2 = TOP + Inches(1.6)

    loan_cards = [
        ("融資の使途", C_ACCENT, [
            "内装工事・設計費：700万円",
            "厨房設備・什器：450万円",
            "保証金・礼金：150万円",
            "備品・食器類：120万円",
            "広告宣伝（開業前）：100万円",
            "予備費：80万円",
        ]),
        ("返済の根拠", RGBColor(0x4C, 0xAF, 0x50), [
            "月商720万円時の営業利益：約75万円",
            "月次借入返済額：6.4万円",
            "返済比率：売上の0.9%",
            "BEP 460万円達成後、即返済可能",
            "手元資金1,030万円がバッファ",
            "累積赤字試算：最大▲240万円",
        ]),
        ("審査強化ポイント", RGBColor(0x42, 0xA5, 0xF5), [
            "自己資本比率 26%（700/2,630万円）",
            "資本金700万円（準資本含む）",
            "チェーン本部長 13年の実務経験",
            "詳細なOPSマニュアル作成済み",
            "競合分析・立地調査完了",
            "事業計画書・月次P&L完備",
        ]),
    ]

    for i, (title, color, lines) in enumerate(loan_cards):
        x = PAD + i * (W3 + GAP)
        card(slide, x, TOP2, W3, Inches(4.5), title, lines, color)

    # 投資家向け追記（最下部）
    section_rect(slide, PAD, Inches(6.85), SW - PAD * 2, Inches(0.5),
                 RGBColor(0x1A, 0x12, 0x20), RGBColor(0x66, 0x44, 0xCC), 0.5)
    tb(slide, "　📌 エンジェル投資家向け：3年目FC展開フェーズに合わせてSeries A 3,000〜5,000万円の調達を予定しております。ご関心の方はご相談ください。",
       PAD + Inches(0.1), Inches(6.9), SW - PAD * 2 - Inches(0.2), Inches(0.4),
       size=9, color=RGBColor(0xBB, 0xAA, 0xFF))


def s_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, 0, Inches(0.07))
    accent_bar(slide, SH - Inches(0.07), Inches(0.07))

    tb(slide, "然", Inches(5.8), Inches(1.0), Inches(1.7), Inches(1.5),
       size=72, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb(slide, "肉酒場 然", Inches(3.5), Inches(2.5), Inches(6.33), Inches(0.8),
       size=32, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, "「糀漬け肉の定食と、和の一杯。」",
       Inches(2.5), Inches(3.4), Inches(8.33), Inches(0.55),
       size=16, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)

    tb(slide, "ご検討のほど、よろしくお願いいたします。",
       Inches(2.5), Inches(4.3), Inches(8.33), Inches(0.55),
       size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 連絡先
    section_rect(slide, Inches(4.0), Inches(5.2), Inches(5.33), Inches(1.5),
                 C_SECTION, C_ACCENT, 0.5)
    tb(slide, "お問い合わせ",
       Inches(4.15), Inches(5.3), Inches(5.0), Inches(0.35),
       size=10, color=C_ACCENT)
    tb(slide, "代表：○○　○○\nhunkotusaisinryo@gmail.com",
       Inches(4.15), Inches(5.65), Inches(5.0), Inches(0.9),
       size=12, color=C_LIGHT)


# ===== メイン =====

def generate(output_dir="reports"):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    # blank layout
    while len(prs.slide_layouts) < 7:
        prs.slide_layouts[0]

    s_cover(prs)
    s_executive_summary(prs)
    s_market(prs)
    s_concept(prs)
    s_pl_monthly(prs)
    s_funding(prs)
    s_growth_plan(prs)
    s_revenue_chart(prs)
    s_risk(prs)
    s_team(prs)
    s_ask(prs)
    s_closing(prs)

    os.makedirs(output_dir, exist_ok=True)
    today = datetime.date.today().strftime("%Y%m%d")
    out = f"{output_dir}/肉酒場然_融資資料_{today}.pptx"
    prs.save(out)
    print(f"生成完了: {out}")
    return out


if __name__ == "__main__":
    generate()
