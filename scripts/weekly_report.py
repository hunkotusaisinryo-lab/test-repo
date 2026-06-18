#!/usr/bin/env python3
"""
週次トレンドリサーチレポート生成スクリプト
毎週月曜日にGitHub Actionsから実行される
前週（月〜日）のresearch/*.mdをPowerPointにまとめる
"""

import os
import re
import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# カラー定義
COLOR_BG = RGBColor(0x1A, 0x1A, 0x2E)        # 濃紺（背景）
COLOR_ACCENT = RGBColor(0xE8, 0x9C, 0x45)     # 金（アクセント）
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # 白
COLOR_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)      # ライトグレー
COLOR_SECTION = RGBColor(0x2D, 0x2D, 0x44)    # セクション背景


def get_last_week_dates(research_dir="research"):
    """前週（月〜日）の日付リストを返す。前週にファイルがなければ直近7日分にフォールバック"""
    today = datetime.date.today()
    last_monday = today - datetime.timedelta(days=today.weekday() + 7)
    last_week = [last_monday + datetime.timedelta(days=i) for i in range(7)]

    # 前週にファイルが存在するか確認
    research_path = Path(research_dir)
    if any((research_path / f"{d}.md").exists() for d in last_week):
        return last_week

    # フォールバック: 直近14日以内の実在ファイルを最大7件取得
    available = sorted([
        datetime.date.fromisoformat(p.stem)
        for p in research_path.glob("????-??-??.md")
        if (today - datetime.date.fromisoformat(p.stem)).days <= 14
    ])
    if available:
        return available[-7:]

    return last_week


def parse_markdown(filepath):
    """markdownファイルをセクションごとに解析する"""
    with open(filepath, encoding="utf-8") as f:
        content = f.read()

    sections = {}
    current_section = None
    current_lines = []

    for line in content.split("\n"):
        if line.startswith("## "):
            if current_section:
                sections[current_section] = "\n".join(current_lines).strip()
            current_section = line[3:].strip()
            current_lines = []
        elif line.startswith("# "):
            sections["title"] = line[2:].strip()
        elif current_section:
            current_lines.append(line)

    if current_section:
        sections[current_section] = "\n".join(current_lines).strip()

    return sections


def extract_bullets(text, max_items=5):
    """テキストから箇条書き項目を抽出する"""
    bullets = []
    for line in text.split("\n"):
        line = line.strip()
        if line.startswith("- ") or line.startswith("・"):
            item = line.lstrip("- ").lstrip("・").strip()
            if item:
                bullets.append(item)
        elif line.startswith("**") and line.endswith("**"):
            bullets.append(line.strip("*"))
    return bullets[:max_items]


def set_slide_background(slide, prs, color):
    """スライド背景色を設定する"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    """テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_title_slide(prs, week_start, week_end, total_days):
    """表紙スライドを作成する"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, COLOR_BG)

    W = prs.slide_width
    H = prs.slide_height

    # アクセントライン（上部）
    from pptx.util import Pt as PtUtil
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    # メインタイトル
    add_textbox(slide, "飲食トレンド週次レポート",
                Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.2),
                font_size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # 期間
    period = f"{week_start.strftime('%Y年%m月%d日')}（月）〜{week_end.strftime('%m月%d日')}（日）"
    add_textbox(slide, period,
                Inches(1.0), Inches(3.5), Inches(11.33), Inches(0.7),
                font_size=20, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

    # サブテキスト
    add_textbox(slide, f"対象日数：{total_days}日分　｜　自動生成レポート",
                Inches(1.0), Inches(4.3), Inches(11.33), Inches(0.5),
                font_size=14, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    # アクセントライン（下部）
    shape2 = slide.shapes.add_shape(
        1, 0, H - Inches(0.08), W, Inches(0.08)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLOR_ACCENT
    shape2.line.fill.background()

    return slide


def add_section_block(slide, title, body_text, left, top, width, height, accent_color, text_color):
    """セクションタイトル＋本文のブロックを描画する"""
    # 背景
    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_SECTION
    bg.line.fill.background()

    # 左アクセントバー
    bar = slide.shapes.add_shape(1, left, top, Inches(0.04), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # タイトル
    add_textbox(slide, title,
                left + Inches(0.12), top + Inches(0.08),
                width - Inches(0.2), Inches(0.3),
                font_size=10, bold=True, color=accent_color)

    # 本文
    add_textbox(slide, body_text,
                left + Inches(0.12), top + Inches(0.4),
                width - Inches(0.2), height - Inches(0.48),
                font_size=9, color=text_color, wrap=True)


def add_daily_slide(prs, date, sections):
    """1日分のスライドを作成する（16:9レイアウト）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, COLOR_BG)

    W = prs.slide_width   # 13.33inch
    H = prs.slide_height  # 7.5inch

    # 上部アクセントバー
    shape = slide.shapes.add_shape(1, 0, 0, W, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    # ヘッダー背景
    header_h = Inches(0.75)
    header_bg = slide.shapes.add_shape(1, 0, Inches(0.05), W, header_h)
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = COLOR_SECTION
    header_bg.line.fill.background()

    weekdays = ["月", "火", "水", "木", "金", "土", "日"]
    weekday = weekdays[date.weekday()]
    date_str = f"{date.strftime('%Y年%m月%d日')}（{weekday}）"
    add_textbox(slide, date_str,
                Inches(0.4), Inches(0.1), Inches(8), Inches(0.6),
                font_size=22, bold=True, color=COLOR_WHITE)

    # ========== レイアウト定数 ==========
    PAD = Inches(0.25)         # 外余白
    GAP = Inches(0.18)         # ブロック間隔
    content_top = Inches(0.05) + header_h + Inches(0.18)
    content_h = H - content_top - Inches(0.15)  # 有効高さ

    # 3カラム構成
    col_w = (W - PAD * 2 - GAP * 2) / 3
    col1_x = PAD
    col2_x = col1_x + col_w + GAP
    col3_x = col2_x + col_w + GAP

    # --- 左カラム: ハイライト（上） + SNSトレンド（下） ---
    highlight_h = content_h * 0.52
    sns_h = content_h - highlight_h - GAP

    highlight_text = sections.get("今日のハイライト", "（データなし）")
    highlight_clean = re.sub(r'\*+', '', highlight_text).strip()
    add_section_block(slide, "今日のハイライト", highlight_clean[:280],
                      col1_x, content_top, col_w, highlight_h,
                      COLOR_ACCENT, COLOR_LIGHT)

    sns_bullets = extract_bullets(sections.get("SNSトレンド", ""), max_items=4)
    sns_text = "\n".join(f"• {b}" for b in sns_bullets) if sns_bullets else "（データなし）"
    add_section_block(slide, "SNSトレンド", sns_text,
                      col1_x, content_top + highlight_h + GAP, col_w, sns_h,
                      RGBColor(0x5B, 0xB8, 0xD4), COLOR_LIGHT)

    # --- 中カラム: 話題の料理・食材（上） + 注目の業態（下） ---
    food_h = content_h * 0.52
    store_h = content_h - food_h - GAP

    food_bullets = extract_bullets(sections.get("話題の料理・食材", ""), max_items=5)
    food_text = "\n".join(f"• {b}" for b in food_bullets) if food_bullets else "（データなし）"
    add_section_block(slide, "話題の料理・食材", food_text,
                      col2_x, content_top, col_w, food_h,
                      RGBColor(0xE8, 0x7C, 0x3C), COLOR_LIGHT)

    store_bullets = extract_bullets(sections.get("注目の業態・新店舗", ""), max_items=4)
    store_text = "\n".join(f"• {b}" for b in store_bullets) if store_bullets else "（データなし）"
    add_section_block(slide, "注目の業態・新店舗", store_text,
                      col2_x, content_top + food_h + GAP, col_w, store_h,
                      RGBColor(0x7C, 0xC8, 0x7C), COLOR_LIGHT)

    # --- 右カラム: 発酵・健康（上） + 事業アイデア（下） ---
    ferment_h = content_h * 0.40
    idea_h = content_h - ferment_h - GAP

    ferment_bullets = extract_bullets(sections.get("発酵・健康・肉料理関連", ""), max_items=4)
    ferment_text = "\n".join(f"• {b}" for b in ferment_bullets) if ferment_bullets else "（データなし）"
    add_section_block(slide, "発酵・健康・肉料理", ferment_text,
                      col3_x, content_top, col_w, ferment_h,
                      RGBColor(0xC8, 0x96, 0xE8), COLOR_LIGHT)

    # 事業アイデアメモ（アクセント枠）
    idea_top = content_top + ferment_h + GAP
    idea_bg = slide.shapes.add_shape(1, col3_x, idea_top, col_w, idea_h)
    idea_bg.fill.solid()
    idea_bg.fill.fore_color.rgb = RGBColor(0x22, 0x18, 0x08)
    idea_bg.line.color.rgb = COLOR_ACCENT
    idea_bg.line.width = Pt(1.5)

    bar3 = slide.shapes.add_shape(1, col3_x, idea_top, Inches(0.04), idea_h)
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = COLOR_ACCENT
    bar3.line.fill.background()

    add_textbox(slide, "💡 事業アイデアメモ",
                col3_x + Inches(0.12), idea_top + Inches(0.08),
                col_w - Inches(0.2), Inches(0.3),
                font_size=10, bold=True, color=COLOR_ACCENT)

    idea_bullets = extract_bullets(sections.get("事業へのアイデアメモ", ""), max_items=3)
    idea_text = "\n".join(f"• {b}" for b in idea_bullets) if idea_bullets else "（データなし）"
    add_textbox(slide, idea_text,
                col3_x + Inches(0.12), idea_top + Inches(0.42),
                col_w - Inches(0.2), idea_h - Inches(0.5),
                font_size=9, color=COLOR_LIGHT, wrap=True)

    return slide


def add_summary_slide(prs, all_ideas, week_start, week_end):
    """週次まとめスライドを作成する"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, COLOR_BG)

    W = prs.slide_width
    H = prs.slide_height

    shape = slide.shapes.add_shape(1, 0, 0, W, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    header_bg = slide.shapes.add_shape(1, 0, Inches(0.05), W, Inches(0.75))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = COLOR_SECTION
    header_bg.line.fill.background()

    add_textbox(slide, "週次まとめ｜事業アイデア総覧",
                Inches(0.4), Inches(0.1), Inches(10), Inches(0.6),
                font_size=22, bold=True, color=COLOR_WHITE)

    period = f"{week_start.strftime('%m/%d')}〜{week_end.strftime('%m/%d')} の気づき"
    add_textbox(slide, period,
                Inches(10.5), Inches(0.18), Inches(2.5), Inches(0.45),
                font_size=11, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT)

    # アイデアカード一覧
    PAD = Inches(0.35)
    GAP = Inches(0.2)
    y = Inches(1.0)
    row_h = Inches(0.72)

    weekdays = ["月", "火", "水", "木", "金", "土", "日"]
    for date, ideas in all_ideas.items():
        if y + row_h > H - Inches(0.2):
            break

        # 日付バッジ背景
        badge = slide.shapes.add_shape(1, PAD, y, Inches(1.6), row_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        badge.line.color.rgb = COLOR_ACCENT
        badge.line.width = Pt(0.75)

        day_label = f"{date.strftime('%m/%d')}\n（{weekdays[date.weekday()]}）"
        add_textbox(slide, day_label,
                    PAD + Inches(0.05), y + Inches(0.06),
                    Inches(1.5), row_h - Inches(0.1),
                    font_size=11, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

        # アイデア本文背景
        idea_x = PAD + Inches(1.6) + Inches(0.12)
        idea_w = W - idea_x - PAD
        idea_bg = slide.shapes.add_shape(1, idea_x, y, idea_w, row_h)
        idea_bg.fill.solid()
        idea_bg.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x38)
        idea_bg.line.fill.background()

        if ideas:
            idea_text = "　•  " + "\n　•  ".join(ideas[:3])
        else:
            idea_text = "　（記録なし）"
        add_textbox(slide, idea_text,
                    idea_x + Inches(0.1), y + Inches(0.08),
                    idea_w - Inches(0.2), row_h - Inches(0.14),
                    font_size=10, color=COLOR_LIGHT)

        y += row_h + GAP

    return slide


def add_business_idea_slide(prs, all_sections, week_start, week_end):
    """事業アイデア深掘りスライド（毎週必ず生成）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, COLOR_BG)

    W = prs.slide_width
    H = prs.slide_height

    # ヘッダー
    shape = slide.shapes.add_shape(1, 0, 0, W, Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    header_bg = slide.shapes.add_shape(1, 0, Inches(0.05), W, Inches(0.75))
    header_bg.fill.solid(); header_bg.fill.fore_color.rgb = COLOR_SECTION
    header_bg.line.fill.background()

    add_textbox(slide, "💡 今週の事業アイデア｜肉酒場 然 への活用",
                Inches(0.4), Inches(0.1), Inches(10), Inches(0.6),
                font_size=22, bold=True, color=COLOR_WHITE)

    period = f"{week_start.strftime('%m/%d')}〜{week_end.strftime('%m/%d')}"
    add_textbox(slide, period,
                Inches(11.0), Inches(0.18), Inches(2.0), Inches(0.45),
                font_size=11, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT)

    # 週のトレンドから自動でアイデアを生成
    PAD = Inches(0.35)
    TOP = Inches(0.95)
    GAP = Inches(0.18)
    W3 = (W - PAD * 2 - GAP * 2) / 3

    # 全日のデータから各セクションをまとめる
    all_highlights = []
    all_foods = []
    all_ferments = []
    all_sns = []
    all_ideas = []

    for sections in all_sections.values():
        highlight = re.sub(r'\*+', '', sections.get("今日のハイライト", "")).strip()
        if highlight:
            all_highlights.append(highlight[:80])
        all_foods += extract_bullets(sections.get("話題の料理・食材", ""), max_items=2)
        all_ferments += extract_bullets(sections.get("発酵・健康・肉料理関連", ""), max_items=2)
        all_sns += extract_bullets(sections.get("SNSトレンド", ""), max_items=2)
        all_ideas += extract_bullets(sections.get("事業へのアイデアメモ", ""), max_items=3)

    # 左カラム：今週の注目トレンド
    col1_x = PAD
    bg1 = slide.shapes.add_shape(1, col1_x, TOP, W3, H - TOP - Inches(0.2))
    bg1.fill.solid(); bg1.fill.fore_color.rgb = COLOR_SECTION
    bg1.line.fill.background()
    bar1 = slide.shapes.add_shape(1, col1_x, TOP, Inches(0.04), H - TOP - Inches(0.2))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = COLOR_ACCENT
    bar1.line.fill.background()

    add_textbox(slide, "今週のキートレンド",
                col1_x + Inches(0.12), TOP + Inches(0.08), W3 - Inches(0.2), Inches(0.3),
                font_size=11, bold=True, color=COLOR_ACCENT)

    trend_lines = []
    for f in all_foods[:4]:
        trend_lines.append(f"🍖 {f[:45]}")
    for s in all_ferments[:3]:
        trend_lines.append(f"🧪 {s[:45]}")
    for sn in all_sns[:2]:
        trend_lines.append(f"📱 {sn[:45]}")

    add_textbox(slide, "\n".join(trend_lines) if trend_lines else "（データなし）",
                col1_x + Inches(0.12), TOP + Inches(0.45), W3 - Inches(0.2), H - TOP - Inches(0.7),
                font_size=9, color=COLOR_LIGHT, wrap=True)

    # 中カラム：日次アイデアまとめ
    col2_x = PAD + W3 + GAP
    bg2 = slide.shapes.add_shape(1, col2_x, TOP, W3, H - TOP - Inches(0.2))
    bg2.fill.solid(); bg2.fill.fore_color.rgb = COLOR_SECTION
    bg2.line.fill.background()
    bar2 = slide.shapes.add_shape(1, col2_x, TOP, Inches(0.04), H - TOP - Inches(0.2))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(0x5B, 0xB8, 0x6C)
    bar2.line.fill.background()

    add_textbox(slide, "日次アイデアメモ（今週分）",
                col2_x + Inches(0.12), TOP + Inches(0.08), W3 - Inches(0.2), Inches(0.3),
                font_size=11, bold=True, color=RGBColor(0x5B, 0xB8, 0x6C))

    idea_lines = [f"• {i[:55]}" for i in all_ideas[:8]] if all_ideas else ["（今週のアイデアメモなし）"]
    add_textbox(slide, "\n".join(idea_lines),
                col2_x + Inches(0.12), TOP + Inches(0.45), W3 - Inches(0.2), H - TOP - Inches(0.7),
                font_size=9, color=COLOR_LIGHT, wrap=True)

    # 右カラム：然への具体的活用アイデア（固定テンプレート＋トレンド反映）
    col3_x = PAD + W3 * 2 + GAP * 2
    bg3 = slide.shapes.add_shape(1, col3_x, TOP, W3, H - TOP - Inches(0.2))
    bg3.fill.solid(); bg3.fill.fore_color.rgb = RGBColor(0x22, 0x18, 0x08)
    bg3.line.color.rgb = COLOR_ACCENT; bg3.line.width = Pt(1.5)
    bar3 = slide.shapes.add_shape(1, col3_x, TOP, Inches(0.04), H - TOP - Inches(0.2))
    bar3.fill.solid(); bar3.fill.fore_color.rgb = COLOR_ACCENT
    bar3.line.fill.background()

    add_textbox(slide, "🔥 肉酒場 然 への活用アクション",
                col3_x + Inches(0.12), TOP + Inches(0.08), W3 - Inches(0.2), Inches(0.3),
                font_size=11, bold=True, color=COLOR_ACCENT)

    # トレンドから自動でアクションを生成
    action_items = []
    if all_foods:
        action_items.append(f"【メニュー】{all_foods[0][:30]}を糀漬けアレンジで検討")
    if all_ferments:
        action_items.append(f"【発酵】{all_ferments[0][:30]}を看板メニューに応用")
    if all_sns:
        action_items.append(f"【SNS】{all_sns[0][:30]}の投稿フォーマットを参考に")
    # 固定アクション
    action_items += [
        "【仕込み】今週のトレンド食材を翌週の漬け込み素材に追加",
        "【集客】週次トレンドをInstagramストーリーズで発信",
        "【FC】話題業態のOPSを分析してマニュアルに反映",
    ]

    action_text = "\n\n".join([f"▶ {a}" for a in action_items[:5]])
    add_textbox(slide, action_text,
                col3_x + Inches(0.12), TOP + Inches(0.45), W3 - Inches(0.2), H - TOP - Inches(0.7),
                font_size=9, color=COLOR_LIGHT, wrap=True)

    return slide


def generate_report(research_dir="research", output_dir="zen/06_週次レポート"):
    dates = get_last_week_dates(research_dir)
    week_start = dates[0]
    week_end = dates[-1]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 表紙
    available_dates = []
    all_sections = {}
    for date in dates:
        filepath = Path(research_dir) / f"{date}.md"
        if filepath.exists():
            available_dates.append(date)
            all_sections[date] = parse_markdown(filepath)

    add_title_slide(prs, week_start, week_end, len(available_dates))

    # 日別スライド
    all_ideas = {}
    for date in available_dates:
        sections = all_sections[date]
        add_daily_slide(prs, date, sections)
        idea_bullets = extract_bullets(sections.get("事業へのアイデアメモ", ""), max_items=2)
        all_ideas[date] = idea_bullets

    # まとめスライド
    if all_ideas:
        add_summary_slide(prs, all_ideas, week_start, week_end)

    # 事業アイデア深掘りスライド（毎週必ず生成）
    add_business_idea_slide(prs, all_sections, week_start, week_end)

    # 保存
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{output_dir}/weekly_report_{week_start}_{week_end}.pptx"
    prs.save(filename)
    print(f"レポート生成完了: {filename}")
    return filename


if __name__ == "__main__":
    generate_report()
