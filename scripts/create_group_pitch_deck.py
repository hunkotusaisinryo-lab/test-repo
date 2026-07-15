#!/usr/bin/env python3
"""
然グループ 統合ピッチデッキ（強化版）生成スクリプト
明日の打ち合わせ用 - ブランドカラー・財務強化・競合優位性対応
"""

import os
import io
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 定数 ───────────────────────────────────────────────
OUTPUT_DIR = Path("/home/user/test-repo/zen/07_融資・投資家資料")
OUTPUT_FILE = OUTPUT_DIR / "然グループ_統合ピッチデッキ_強化版.pptx"
LOGO_DIR    = Path("/home/user/test-repo/zen/08_ロゴ・ブランド資産")
LOGO_FILE   = LOGO_DIR / "zen_brand_logo.png"
LOGO_GROUP  = LOGO_DIR / "ZEN_GROUP_LOGO.png"
LOGO_DASHI  = LOGO_DIR / "ZEN_DASHI_LOGO.png"
LOGO_NIKU   = LOGO_DIR / "ZEN_NIKU_LOGO.png"
LOGO_SOBA   = LOGO_DIR / "ZEN_SOBA_LOGO.png"

FONT_NAME = "IPAGothic"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# グループカラー（ロゴ準拠）
BG_COLOR   = RGBColor(0xFA, 0xFA, 0xF6)
ACCENT     = RGBColor(0x2C, 0x4A, 0x3E)   # 深緑（グループ共通）
GOLD       = RGBColor(0xC4, 0x9A, 0x22)   # ゴールド（出汁然・ロゴ準拠）
DARK_RED   = RGBColor(0x8B, 0x20, 0x20)   # ダークレッド（肉酒場然）
NAVY       = RGBColor(0x1A, 0x2A, 0x4A)   # ネイビー（蕎麦然）
TEXT_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF0, 0xEE, 0xE8)

# matplotlib用カラー
MPL_GREEN   = "#2C4A3E"
MPL_GOLD    = "#C49A22"
MPL_RED     = "#8B2020"
MPL_NAVY    = "#1A2A4A"
MPL_LIGHT   = "#F0EEE8"

HEADER_H = Inches(0.95)

# ─── ヘルパー関数 ────────────────────────────────────────

def set_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color if color else BG_COLOR


def add_rect(slide, l, t, w, h, color, line=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def add_header_band(slide, title, color=None, show_logo=True):
    c = color if color else ACCENT
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, c)
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(10.5), HEADER_H - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = Pt(26)
    run.font.color.rgb = WHITE
    run.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    # ロゴを右上に配置
    if show_logo and LOGO_FILE.exists():
        slide.shapes.add_picture(str(LOGO_FILE), Inches(10.8), Inches(0.05), Inches(2.3), Inches(0.82))


def txt(slide, text, l, t, w, h, size=14, color=None, bold=False,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    if color is None:
        color = TEXT_COLOR
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def multi_txt(slide, lines, l, t, w, h, size=13, color=None,
              bold=False, align=PP_ALIGN.LEFT, spacing=None):
    if color is None:
        color = TEXT_COLOR
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def table(slide, headers, rows, l, t, w, h, hdr_color=None, font_size=11):
    if hdr_color is None:
        hdr_color = ACCENT
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = h // n_rows
    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    col_w = w // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        bg = LIGHT_BG if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_COLOR


def embed_figure(slide, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, l, t, w, h)
    plt.close(fig)


def kpi_box(slide, value, label, l, t, w, h, color=None):
    c = color if color else ACCENT
    add_rect(slide, l, t, w, h, c)
    txt(slide, value, l, t + Inches(0.08), w, Inches(0.65),
        size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + Inches(0.68), w, Inches(0.38),
        size=10, color=WHITE, align=PP_ALIGN.CENTER)


def section_box(slide, title, items, l, t, w, title_color=None):
    c = title_color if title_color else ACCENT
    add_rect(slide, l, t, w, Inches(0.38), c)
    txt(slide, title, l + Inches(0.1), t + Inches(0.04), w - Inches(0.2), Inches(0.32),
        size=12, bold=True, color=WHITE)
    y = t + Inches(0.42)
    for item in items:
        txt(slide, "▸ " + item, l + Inches(0.1), y, w - Inches(0.2), Inches(0.38),
            size=11, color=TEXT_COLOR)
        y += Inches(0.38)


def blank(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, ACCENT if dark else None)
    return slide


# ─── グラフ生成 ───────────────────────────────────────────

def fig_growth_chart():
    """5カ年 月商推移グラフ"""
    fig, ax = plt.subplots(figsize=(7, 3.8))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    years = ["Y1\n2026", "Y2\n2027", "Y3\n2028", "Y4\n2029", "Y5\n2030"]
    niku  = [300, 500,  600,  800,  1000]
    dashi = [600, 1500, 3000, 5000, 7000]
    soba  = [300, 1000, 2400, 4200, 6000]
    sns   = [0,   0,    168,  500,  1000]

    x = np.arange(len(years))
    w = 0.18
    ax.bar(x - 1.5*w, niku,  w, label="肉酒場 然", color=MPL_RED,   alpha=0.9)
    ax.bar(x - 0.5*w, dashi, w, label="だし 然",   color=MPL_GOLD,  alpha=0.9)
    ax.bar(x + 0.5*w, soba,  w, label="蕎麦 然",   color=MPL_NAVY,  alpha=0.9)
    ax.bar(x + 1.5*w, sns,   w, label="SNS事業部", color=MPL_GREEN, alpha=0.9)

    totals = [n+d+s+sn for n,d,s,sn in zip(niku,dashi,soba,sns)]
    for xi, tot in zip(x, totals):
        ax.text(xi, tot + 80, f"{tot:,}万", ha="center", va="bottom",
                fontsize=8, color=MPL_GREEN, fontweight="bold",
                fontfamily="IPAGothic")

    ax.set_xticks(x)
    ax.set_xticklabels(years, fontfamily="IPAGothic", fontsize=10)
    ax.set_ylabel("月商（万円）", fontfamily="IPAGothic", fontsize=9)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
    ax.legend(prop={"family": "IPAGothic", "size": 9}, loc="upper left")
    ax.spines[["top","right"]].set_visible(False)
    ax.set_title("グループ月商推移（5カ年）", fontfamily="IPAGothic",
                 fontsize=12, color=MPL_GREEN, fontweight="bold", pad=10)
    fig.subplots_adjust(top=0.88, bottom=0.12, left=0.1, right=0.97)
    return fig


def fig_pl_chart():
    """簡易PL棒グラフ（だし 然 1店舗モデル・利益率15%）"""
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    labels = ["売上", "原価", "人件費", "家賃", "その他", "営業利益"]
    values = [960, -317, -270, -60, -169, 144]
    colors = [MPL_GOLD if v > 0 else ("#C0392B" if i < 5 else MPL_GREEN)
              for i, v in enumerate(values)]
    colors[-1] = MPL_GREEN

    x_pos = range(len(labels))
    bars = ax.bar(x_pos, [abs(v) for v in values], color=colors, alpha=0.9, width=0.6)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 5,
                f"{val:+,}万", ha="center", va="bottom",
                fontsize=9, fontfamily="IPAGothic", fontweight="bold",
                color=MPL_GREEN if val > 0 else "#C0392B")

    ax.set_ylabel("金額（万円/月）", fontfamily="IPAGothic", fontsize=9)
    ax.set_title("だし 然 月次P&L（1店舗・利益率15%）", fontfamily="IPAGothic",
                 fontsize=11, color=MPL_GOLD, fontweight="bold")
    ax.spines[["top","right"]].set_visible(False)
    ax.set_xticks(list(x_pos))
    ax.set_xticklabels(labels, fontfamily="IPAGothic", fontsize=9)
    fig.tight_layout()
    return fig


def fig_group_pl():
    """グループ各ブランド 月次PL比較（1店舗）"""
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    brands = ["だし 然", "肉酒場 然", "蕎麦 然"]
    sales   = [960, 1200, 350]
    profits = [144, 204,   70]
    margins = [15,  17,   20]
    colors  = [MPL_GOLD, MPL_RED, MPL_NAVY]

    x = np.arange(len(brands))
    w = 0.35
    ax.bar(x - w/2, sales,   w, label="月次売上",   color=colors, alpha=0.4)
    bars2 = ax.bar(x + w/2, profits, w, label="月次営業利益", color=colors, alpha=0.9)

    for bar, val, pct in zip(bars2, profits, margins):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 6,
                f"{val}万\n({pct:.0f}%)", ha="center", va="bottom",
                fontsize=8, fontfamily="IPAGothic", fontweight="bold",
                color=MPL_GREEN)

    ax.set_xticks(x)
    ax.set_xticklabels(brands, fontfamily="IPAGothic", fontsize=10)
    ax.set_ylabel("金額（万円/月）", fontfamily="IPAGothic", fontsize=9)
    ax.set_title("各ブランド 月次PL比較（1店舗）", fontfamily="IPAGothic",
                 fontsize=11, color=MPL_GREEN, fontweight="bold")
    ax.legend(prop={"family": "IPAGothic", "size": 9})
    ax.spines[["top","right"]].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
    fig.tight_layout()
    return fig


def fig_recovery_chart():
    """投資回収シミュレーション"""
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    months = list(range(0, 31))
    dashi_cum = []
    soba_cum = []
    init_dashi = -3000
    init_soba = -930

    for m in months:
        if m == 0:
            dashi_cum.append(init_dashi)
            soba_cum.append(init_soba)
        else:
            profit_d = 60 if m <= 3 else (100 if m <= 6 else (130 if m <= 12 else 144))
            profit_s = 25 if m <= 3 else (45 if m <= 6 else 60)
            dashi_cum.append(dashi_cum[-1] + profit_d)
            soba_cum.append(soba_cum[-1] + profit_s)

    ax.plot(months, dashi_cum, color=MPL_GOLD, linewidth=2.5, label="だし 然（初期投資3,000万）", marker="o", markersize=3)
    ax.plot(months, soba_cum, color=MPL_NAVY, linewidth=2.5, label="蕎麦 然（初期投資930万）", marker="s", markersize=3)
    ax.axhline(0, color=MPL_RED, linestyle="--", linewidth=1.5, alpha=0.7)
    ax.fill_between(months, 0, [max(0, v) for v in dashi_cum], alpha=0.1, color=MPL_GOLD)
    ax.fill_between(months, 0, [max(0, v) for v in soba_cum], alpha=0.1, color=MPL_NAVY)

    ax.set_xlabel("経過月数", fontfamily="IPAGothic", fontsize=9)
    ax.set_ylabel("累計損益（万円）", fontfamily="IPAGothic", fontsize=9)
    ax.set_title("投資回収シミュレーション", fontfamily="IPAGothic",
                 fontsize=11, color=MPL_GREEN, fontweight="bold")
    ax.legend(prop={"family": "IPAGothic", "size": 8})
    ax.spines[["top","right"]].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
    fig.tight_layout()
    return fig


def fig_competitive_map():
    """競合ポジショニングマップ"""
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    competitors = [
        # 然グループ3ブランド
        ("だし 然", 4.0, 3.5, MPL_GOLD, 130, True),
        ("肉酒場 然", 4.5, 4.2, MPL_RED, 130, True),
        ("蕎麦 然", 2.5, 4.8, MPL_NAVY, 130, True),
        # 実在競合（恵比寿・新橋・学芸大学エリア）
        ("みなとや（大手町）", 3.5, 2.0, "#888", 80, False),
        ("ひまり堂（恵比寿）", 4.2, 1.8, "#888", 80, False),
        ("ひまり商店（新橋）", 3.0, 2.5, "#888", 80, False),
        ("びゃく（学芸大学）", 4.8, 1.2, "#888", 80, False),
        ("三谷（学芸大学）", 4.6, 1.5, "#888", 80, False),
        ("なかよし（恵比寿）", 3.2, 2.8, "#888", 80, False),
    ]

    for name, x, y, color, size, bold in competitors:
        ax.scatter(x, y, s=size, color=color, alpha=0.85, zorder=3)
        ax.annotate(name, (x, y), textcoords="offset points",
                    xytext=(8, 4), fontfamily="IPAGothic", fontsize=8,
                    fontweight="bold" if bold else "normal",
                    color=color if bold else "gray")

    ax.set_xlim(0, 5.5)
    ax.set_ylim(0, 5.5)
    ax.set_xlabel("素材・体験へのこだわり →", fontfamily="IPAGothic", fontsize=9)
    ax.set_ylabel("FC展開・スケーラビリティ →", fontfamily="IPAGothic", fontsize=9)
    ax.set_title("競合ポジショニングマップ", fontfamily="IPAGothic",
                 fontsize=11, color=MPL_GREEN, fontweight="bold")
    ax.axvline(2.75, color="lightgray", linestyle="--", alpha=0.5)
    ax.axhline(2.75, color="lightgray", linestyle="--", alpha=0.5)
    ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout()
    return fig


def fig_royalty_scale():
    """FCロイヤルティ収入スケール"""
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    fc_count = [5, 10, 15, 20, 25, 30, 40]
    dashi_r  = [c * 960 * 0.03 for c in fc_count]
    soba_r   = [c * 348 * 0.04 for c in fc_count]
    total_r  = [d + s for d, s in zip(dashi_r, soba_r)]

    ax.fill_between(fc_count, dashi_r, alpha=0.6, color=MPL_GOLD, label="だし 然 ロイヤルティ")
    ax.fill_between(fc_count, soba_r, alpha=0.6, color=MPL_NAVY, label="蕎麦 然 ロイヤルティ")
    ax.plot(fc_count, total_r, color=MPL_GREEN, linewidth=2.5, marker="o", markersize=4, label="合計")

    for x, y in zip(fc_count, total_r):
        if x in [10, 20, 40]:
            ax.annotate(f"{int(y):,}万", (x, y), textcoords="offset points",
                        xytext=(4, 6), fontsize=8, fontfamily="IPAGothic",
                        color=MPL_GREEN, fontweight="bold")

    ax.set_xlabel("FC店舗数", fontfamily="IPAGothic", fontsize=9)
    ax.set_ylabel("月次ロイヤルティ収入（万円）", fontfamily="IPAGothic", fontsize=9)
    ax.set_title("FCスケールによるロイヤルティ収入", fontfamily="IPAGothic",
                 fontsize=10, color=MPL_GREEN, fontweight="bold")
    ax.legend(prop={"family": "IPAGothic", "size": 8})
    ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout()
    return fig


def fig_annual_pl_chart():
    """グループ年次売上・EBITDA推移（5カ年）"""
    fig, ax1 = plt.subplots(figsize=(7, 3.8))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax1.set_facecolor(MPL_LIGHT)

    years  = ["Y1\n2026", "Y2\n2027", "Y3\n2028", "Y4\n2029", "Y5\n2030"]
    sales  = [14000, 36000, 72000, 120000, 200000]
    ebitda = [2260,  7140,  14880, 27800,  40000]
    margin = [e/s*100 for e, s in zip(ebitda, sales)]

    x = np.arange(len(years))
    ax1.bar(x, sales, 0.5, label="売上高", color=MPL_GOLD, alpha=0.7)
    ax1.bar(x, ebitda, 0.5, label="EBITDA", color=MPL_GREEN, alpha=0.9)

    ax2 = ax1.twinx()
    ax2.plot(x, margin, color=MPL_RED, linewidth=2.5, marker="o", markersize=5,
             label="EBITDA率")
    for xi, m in zip(x, margin):
        ax2.annotate(f"{m:.0f}%", (xi, m), textcoords="offset points",
                     xytext=(0, 8), fontsize=8, ha="center",
                     fontfamily="IPAGothic", color=MPL_RED, fontweight="bold")

    ax1.set_xticks(x)
    ax1.set_xticklabels(years, fontfamily="IPAGothic", fontsize=9)
    ax1.set_ylabel("金額（万円）", fontfamily="IPAGothic", fontsize=9)
    ax2.set_ylabel("EBITDA率（%）", fontfamily="IPAGothic", fontsize=9, color=MPL_RED)
    ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, prop={"family": "IPAGothic", "size": 8},
               loc="upper left")
    ax1.spines[["top","right"]].set_visible(False)
    ax1.set_title("グループ年次売上・EBITDA推移", fontfamily="IPAGothic",
                  fontsize=11, color=MPL_GREEN, fontweight="bold")
    fig.tight_layout()
    return fig


def fig_cashflow_chart():
    """グループ年次キャッシュフロー"""
    fig, ax = plt.subplots(figsize=(6.5, 3.5))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    years  = ["調達時", "Y1末", "Y2末", "Y3末", "Y4末", "Y5末"]
    op_cf  = [0, 2600, 7500, 15000, 28000, 48000]
    inv_cf = [-3800, -1000, -5000, -8000, -10000, -5000]
    net    = [o + i for o, i in zip(op_cf, inv_cf)]
    cum    = []
    s = 3800
    for n in net:
        s += n
        cum.append(s)

    x = np.arange(len(years))
    ax.bar(x, [max(0, n) for n in net], 0.5, color=MPL_GREEN, alpha=0.8, label="純CF（年次）")
    ax.bar(x, [min(0, n) for n in net], 0.5, color=MPL_RED, alpha=0.7)
    ax.plot(x, cum, color=MPL_GOLD, linewidth=2.5, marker="o", markersize=5, label="累計CF（資金調達後）")
    ax.axhline(0, color="gray", linestyle="--", linewidth=1)

    for xi, c in zip(x, cum):
        if xi > 0:
            ax.annotate(f"{int(c):,}万", (xi, c), textcoords="offset points",
                        xytext=(0, 8), fontsize=7, ha="center",
                        fontfamily="IPAGothic", color=MPL_GOLD, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(years, fontfamily="IPAGothic", fontsize=9)
    ax.set_ylabel("金額（万円）", fontfamily="IPAGothic", fontsize=9)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v):,}"))
    ax.legend(prop={"family": "IPAGothic", "size": 8})
    ax.spines[["top","right"]].set_visible(False)
    ax.set_title("キャッシュフロー推移（資金調達3,800万起点）", fontfamily="IPAGothic",
                 fontsize=10, color=MPL_GREEN, fontweight="bold")
    fig.tight_layout()
    return fig


# ─── スライド生成 ─────────────────────────────────────────

def slide_01_cover(prs):
    slide = blank(prs)
    # 左帯
    add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, ACCENT)
    # グループロゴ（左帯中央に配置）
    if LOGO_GROUP.exists():
        slide.shapes.add_picture(str(LOGO_GROUP), Inches(0.1), Inches(1.8), Inches(4.3), Inches(1.3))
    elif LOGO_FILE.exists():
        slide.shapes.add_picture(str(LOGO_FILE), Inches(0.25), Inches(1.2), Inches(4.0), Inches(4.0))
    else:
        txt(slide, "然", Inches(0), Inches(0.5), Inches(4.5), Inches(4.5),
            size=200, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # 3ブランドロゴを縦に並べる
    for i, logo in enumerate([LOGO_DASHI, LOGO_NIKU, LOGO_SOBA]):
        if logo.exists():
            slide.shapes.add_picture(str(logo), Inches(0.3), Inches(3.3 + i * 1.1), Inches(3.9), Inches(0.95))
    txt(slide, "ZEN BRAND", Inches(0), Inches(6.55), Inches(4.5), Inches(0.4),
        size=11, color=LIGHT_BG, align=PP_ALIGN.CENTER, italic=True)
    # 右側
    txt(slide, "然グループ", Inches(5), Inches(1.2), Inches(7.8), Inches(1.2),
        size=46, bold=True, color=ACCENT)
    txt(slide, "事業説明資料", Inches(5), Inches(2.4), Inches(7.8), Inches(0.9),
        size=32, bold=True, color=TEXT_COLOR)
    add_rect(slide, Inches(5), Inches(3.4), Inches(6.5), Inches(0.04), GOLD)
    txt(slide, "食の全シーンに、然を。",
        Inches(5), Inches(3.6), Inches(7.8), Inches(0.6),
        size=18, color=GOLD, italic=True)
    multi_txt(slide,
        ["然グループ 事業企画室", "2026年7月"],
        Inches(5), Inches(4.5), Inches(7.8), Inches(0.8),
        size=14, color=TEXT_COLOR)
    # ブランド3色ドット
    for i, c in enumerate([GOLD, DARK_RED, NAVY]):
        add_rect(slide, Inches(5 + i*0.5), Inches(5.6), Inches(0.35), Inches(0.35), c)
    txt(slide, "だし然  ／  肉酒場然  ／  蕎麦然",
        Inches(6.7), Inches(5.6), Inches(5), Inches(0.4),
        size=12, color=TEXT_COLOR)
    # フッター
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), ACCENT)
    txt(slide, "Confidential  ─  然グループ Internal Document  ─  無断転載禁止",
        Inches(0.3), SLIDE_H - Inches(0.32), Inches(12), Inches(0.28),
        size=9, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_vision(prs):
    slide = blank(prs)
    add_header_band(slide, "グループビジョン")
    # ロゴを右上に小さく配置
    if LOGO_FILE.exists():
        slide.shapes.add_picture(str(LOGO_FILE), Inches(10.8), Inches(0.05), Inches(2.3), Inches(0.82))
    txt(slide, "「食の全シーンに、然を。」",
        Inches(1), Inches(1.1), Inches(11), Inches(0.9),
        size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, "日本の「だし・発酵・素材」文化を現代の食卓とビジネスに再定義する",
        Inches(1.5), Inches(2.0), Inches(10), Inches(0.5),
        size=15, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    brands = [
        ("然 出汁", "ZEN DASHI", "だしダイニング",
         "出汁体験×インバウンド\nランチ定食＋体験ディナー\n客単価 1,800〜9,000円", GOLD),
        ("然 肉と発酵", "ZEN NIKU", "炭火と発酵の居酒屋",
         "発酵×炭火グリル\n本格居酒屋体験\n客単価 4,000〜6,000円", DARK_RED),
        ("然 立ちそば", "ZEN SOBA", "現代的な立ち食いそば",
         "2分提供・FC展開特化\n駅前・オフィス街\n客単価 1,500円", NAVY),
    ]
    for i, (name, eng, desc, body, color) in enumerate(brands):
        lx = Inches(0.4 + i * 4.3)
        add_rect(slide, lx, Inches(2.8), Inches(4.0), Inches(4.3), color)
        txt(slide, name, lx, Inches(2.9), Inches(4.0), Inches(0.7),
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, eng, lx, Inches(3.55), Inches(4.0), Inches(0.35),
            size=10, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
        add_rect(slide, lx + Inches(0.3), Inches(3.95), Inches(3.4), Inches(0.03), WHITE)
        txt(slide, desc, lx + Inches(0.1), Inches(4.05), Inches(3.8), Inches(0.35),
            size=10, color=WHITE, align=PP_ALIGN.CENTER)
        multi_txt(slide, body.split("\n"), lx + Inches(0.2), Inches(4.5),
                  Inches(3.6), Inches(2.2), size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide_03_market(prs):
    slide = blank(prs)
    add_header_band(slide, "市場機会  ─  なぜ今か")

    kpis = [
        ("25兆円", "国内外食産業\n市場規模", ACCENT),
        ("4.3兆円", "インバウンド消費\n（2025年実績）", GOLD),
        ("3,188万人", "訪日外国人数\n（過去最高水準）", DARK_RED),
        ("68%", "健康志向食品市場\n過去5年成長率", NAVY),
    ]
    for i, (val, label, color) in enumerate(kpis):
        lx = Inches(0.3 + i * 3.2)
        kpi_box(slide, val, label, lx, Inches(1.1), Inches(2.9), Inches(1.3), color)

    left_items = [
        "外食産業はコロナ後の回復を超え、2030年に向けて成長軌道",
        "インバウンド旅行者の「本物の日本食体験」需要が急拡大",
        "発酵食品・出汁・和食の健康志向は国内外で加速中",
        "立ち食い・スタンド業態：コスパ・時短ニーズで拡大継続",
    ]
    right_items = [
        "飲食FC市場：初期投資の軽さから加盟希望者が増加傾向",
        "SNSによる集客格差が顕在化（勝者総取り化が進行）",
        "円安継続→外国人の日本食消費単価は上昇トレンド",
        "2030年大阪万博：関西インバウンド需要の更なる拡大",
    ]
    section_box(slide, "追い風トレンド（国内）", left_items,
                Inches(0.3), Inches(2.7), Inches(6.2), ACCENT)
    section_box(slide, "追い風トレンド（市場環境）", right_items,
                Inches(6.8), Inches(2.7), Inches(6.2), GOLD)

    txt(slide, "→ だし・発酵・素材こだわり × FC展開 × インバウンド対応 ── この3軸が同時に刺さる市場環境",
        Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.6),
        size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_04_overview(prs):
    slide = blank(prs)
    add_header_band(slide, "然グループ  全体像")
    headers = ["ブランド", "業態", "月次売上/店", "利益率", "家賃/月", "月次利益/店"]
    rows = [
        ["然 出汁", "出汁定食・体験ディナー", "960万円", "15%", "60万円（25坪）", "144万円"],
        ["然 肉と発酵", "炭火コース＋宴会居酒屋", "1,200万円", "17%", "80万円（38坪）", "204万円"],
        ["然 立ちそば", "立ち食い蕎麦（8坪・小型）", "350万円", "20%", "45万円（8坪）", "70万円"],
        ["SNS事業部", "飲食特化SNSコンサル", "月9.8〜39.8万円", "─", "─", "83万円"],
    ]
    table(slide, headers, rows, Inches(0.3), Inches(1.1), Inches(12.7), Inches(2.5))

    txt(slide, "食の時間帯フルカバー戦略",
        Inches(0.3), Inches(3.85), Inches(4), Inches(0.4),
        size=12, bold=True, color=ACCENT)

    timeline = [
        ("蕎麦 然", "11:00〜15:00 ランチ", NAVY),
        ("だし 然", "11:30〜22:00 ランチ〜ディナー", GOLD),
        ("肉酒場 然", "17:00〜23:00 ディナー", DARK_RED),
    ]
    for i, (name, time, color) in enumerate(timeline):
        lx = Inches(0.3 + i * 4.3)
        add_rect(slide, lx, Inches(4.35), Inches(4.0), Inches(0.55), color)
        txt(slide, name, lx, Inches(4.38), Inches(4.0), Inches(0.32),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, time, lx, Inches(4.7), Inches(4.0), Inches(0.25),
            size=9, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, "→ 全時間帯をカバーし、顧客をグループ内で回遊させる「然エコシステム」",
        Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.4),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # シナジー簡略
    synergies = [
        ("出汁共通仕入", "原価率1〜2%改善"),
        ("然カード会員", "グループ内回遊促進"),
        ("SNS内製", "広告費年300万削減"),
        ("FCクロスセル", "SNS事業受注創出"),
    ]
    for i, (title, body) in enumerate(synergies):
        lx = Inches(0.3 + i * 3.2)
        add_rect(slide, lx, Inches(5.65), Inches(3.0), Inches(1.5), ACCENT)
        txt(slide, title, lx, Inches(5.72), Inches(3.0), Inches(0.45),
            size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, body, lx, Inches(6.2), Inches(3.0), Inches(0.45),
            size=10, color=WHITE, align=PP_ALIGN.CENTER)


def slide_05_dashi(prs):
    slide = blank(prs)
    add_header_band(slide, "だし 然  ─  出汁体験レストラン", GOLD)
    txt(slide, "「一杯の出汁が、今日を整える。」",
        Inches(0.5), Inches(1.1), Inches(12), Inches(0.55),
        size=20, bold=True, color=GOLD)

    section_box(slide, "ビジネスモデル（昼夜二毛作）",
                ["昼：出汁定食 1,600〜1,900円 ／ 地元ワーカー向け 平日22日安定集客",
                 "夜：出汁しゃぶ・体験コース 8,000〜12,000円 ／ インバウンド向け体験型",
                 "昆布×本枯れ節×麹発酵のオリジナルブレンド出汁が核心差別化",
                 "25坪・34席（カウンター10席）・4ヶ国語対応・Alipay/WeChatPay導入",
                 "【家賃】60万円/月（25坪・坪2.4万）│ 原価33% │ 人件費28%"],
                Inches(0.3), Inches(1.8), Inches(6.2), GOLD)

    section_box(slide, "展開計画",
                ["Phase1：東京直営 浅草・銀座・上野",
                 "Phase2：関西FC 京都祇園・大阪道頓堀",
                 "Phase3：20店舗（直営5＋FC15）",
                 "Phase4：海外 台湾・シンガポール（DASHI ZEN）"],
                Inches(6.8), Inches(1.8), Inches(6.2), ACCENT)

    # だし然ロゴ
    if LOGO_DASHI.exists():
        slide.shapes.add_picture(str(LOGO_DASHI), Inches(6.6), Inches(1.75), Inches(2.8), Inches(0.7))
    # P&Lグラフ
    fig = fig_pl_chart()
    embed_figure(slide, fig, Inches(0.3), Inches(3.8), Inches(5.8), Inches(3.3))

    # KPI
    kpis = [("月次売上", "960万円"), ("営業利益", "144万円"), ("利益率", "15%"), ("投資回収", "18〜24ヶ月")]
    for i, (label, val) in enumerate(kpis):
        lx = Inches(6.4 + i * 1.7)
        add_rect(slide, lx, Inches(3.8), Inches(1.5), Inches(1.0), ACCENT)
        txt(slide, val, lx, Inches(3.85), Inches(1.5), Inches(0.55),
            size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, label, lx, Inches(4.42), Inches(1.5), Inches(0.35),
            size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # FC収益
    add_rect(slide, Inches(6.4), Inches(5.0), Inches(6.6), Inches(2.1), LIGHT_BG)
    txt(slide, "FCモデル収益", Inches(6.5), Inches(5.05), Inches(6.4), Inches(0.4),
        size=12, bold=True, color=ACCENT)
    multi_txt(slide,
        ["加盟金：300万円 ／ ロイヤルティ：売上の3%",
         "FC初期投資：3,200万円（標準）",
         "FC月次営業利益：約280万円",
         "FC投資回収期間：約20ヶ月",
         "本部ロイヤルティ収入：FC10店で月288万円"],
        Inches(6.5), Inches(5.5), Inches(6.4), Inches(1.5),
        size=11, color=TEXT_COLOR)


def slide_06_niku(prs):
    slide = blank(prs)
    add_header_band(slide, "然 肉と発酵  ─  炭火と発酵の居酒屋", DARK_RED)
    txt(slide, "「発酵の香りと肉煙が満ちる、大人の隠れ酒場」",
        Inches(0.5), Inches(1.1), Inches(12), Inches(0.55),
        size=20, bold=True, color=DARK_RED)

    left_items = [
        "38坪・55席（個室・半個室15席含む）で宴会需要を獲得",
        "発酵肉（麹漬け・味噌漬け）を炭火グリルで「見せる」体験型設計",
        "ディナーコース主体（8,000〜10,000円）＋平日ランチ炭火定食2,000円",
        "日本酒・焼酎の発酵ペアリングでドリンク単価を底上げ",
        "【家賃】80万円/月（38坪・坪2.1万）│ 原価35% │ 人件費27.5%",
    ]
    right_items = [
        "然グループの「顔」かつ 発酵文化の発信基地",
        "炭火・発酵器具がSNS映えコンテンツになる",
        "常連顧客育成で口コミによる自然集客を実現",
        "SNS事業部との連動でブランド発信を内製化",
    ]
    section_box(slide, "コンセプト・差別化", left_items,
                Inches(0.3), Inches(1.8), Inches(6.2), DARK_RED)
    section_box(slide, "戦略的位置づけ", right_items,
                Inches(6.8), Inches(1.8), Inches(6.2), ACCENT)

    # 肉酒場然ロゴ
    if LOGO_NIKU.exists():
        slide.shapes.add_picture(str(LOGO_NIKU), Inches(6.6), Inches(1.75), Inches(2.8), Inches(0.7))
    kpis = [("月次売上目標", "1,200万円", DARK_RED),
            ("月次営業利益", "200万円", ACCENT),
            ("利益率", "17%", DARK_RED),
            ("客単価", "8,000〜10,000円", ACCENT)]
    for i, (label, val, color) in enumerate(kpis):
        lx = Inches(0.3 + i * 3.2)
        add_rect(slide, lx, Inches(4.7), Inches(3.0), Inches(1.0), color)
        txt(slide, val, lx, Inches(4.75), Inches(3.0), Inches(0.55),
            size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, label, lx, Inches(5.35), Inches(3.0), Inches(0.32),
            size=10, color=WHITE, align=PP_ALIGN.CENTER)

    section_box(slide, "今後の展開",
                ["2号店：東京都内（恵比寿・中目黒エリア）",
                 "フラッグシップ店として然グループのブランド旗艦に",
                 "将来的に大阪・福岡への直営拡大"],
                Inches(0.3), Inches(5.95), Inches(12.7), ACCENT)


def slide_07_soba(prs):
    slide = blank(prs)
    add_header_band(slide, "然 立ちそば  ─  現代的な立ち食いそばFC", NAVY)
    txt(slide, "「2分で届く、本物の蕎麦。」",
        Inches(0.5), Inches(1.1), Inches(12), Inches(0.55),
        size=20, bold=True, color=NAVY)

    left_items = [
        "8坪・8スタンド・1〜2人オペレーション（小型化）",
        "客単価1,500円・回転時間12分・1日80〜90人集客",
        "揚げたて天ぷら・北海道幌加内産蕎麦を使用",
        "QR前払い・キャッシュレス100%で人件費最適化",
        "【家賃】45万円/月（8坪・坪5.6万）│ 原価30% │ 人件費23%",
    ]
    right_items = [
        "FC加盟金200万円・ロイヤルティ4%",
        "標準内装パッケージで施工期間2ヶ月以内",
        "首都圏駅前・オフィス街から展開",
        "小型化でFC初期投資を抑え加盟ハードルを下げる",
        "目標：30店舗（直営3＋FC27）",
    ]
    section_box(slide, "業態設計", left_items,
                Inches(0.3), Inches(1.8), Inches(6.2), NAVY)
    section_box(slide, "FC展開モデル", right_items,
                Inches(6.8), Inches(1.8), Inches(6.2), ACCENT)

    # 蕎麦然ロゴ
    if LOGO_SOBA.exists():
        slide.shapes.add_picture(str(LOGO_SOBA), Inches(6.6), Inches(1.75), Inches(2.8), Inches(0.7))
    # 回収シミュレーショングラフ
    fig = fig_recovery_chart()
    embed_figure(slide, fig, Inches(0.3), Inches(4.0), Inches(6.0), Inches(3.1))

    kpis = [("月次売上", "350万円", NAVY),
            ("月次利益", "70万円", ACCENT),
            ("利益率", "20%", NAVY),
            ("投資回収", "12〜15ヶ月", ACCENT)]
    for i, (label, val, color) in enumerate(kpis):
        lx = Inches(6.5 + i * 1.7)
        add_rect(slide, lx, Inches(4.0), Inches(1.5), Inches(0.95), color)
        txt(slide, val, lx, Inches(4.05), Inches(1.5), Inches(0.52),
            size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, label, lx, Inches(4.6), Inches(1.5), Inches(0.32),
            size=9, color=WHITE, align=PP_ALIGN.CENTER)

    multi_txt(slide,
        ["■ FC収益シミュレーション（FC1店舗・小型8坪）",
         "初期投資合計：800万円（内装・設備・物件取得費）",
         "月次売上：350万円",
         "月次営業利益：約56万円（ロイヤルティ控除後）",
         "投資回収期間：約14〜16ヶ月",
         "本部ロイヤルティ：FC10店で月140万円（売上の4%）"],
        Inches(6.5), Inches(5.1), Inches(6.5), Inches(2.0),
        size=11, color=TEXT_COLOR)


def slide_sns_marketing(prs):
    slide = blank(prs)
    add_header_band(slide, "販売・SNSマーケティング戦略")
    txt(slide, "「SNSで集客する飲食グループ」を自ら実証し、外販へ",
        Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
        size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # フェーズ帯
    phases = [
        ("PHASE 1\n開業前3〜6ヶ月", "外部SNSコンサル起用\nスタートダッシュ", DARK_RED),
        ("PHASE 2\nY1〜Y2", "内製移行＋毎日トレンドリサーチ\nコンテンツ量産体制", GOLD),
        ("PHASE 3\nY3〜", "SNS事業部外販スタート\n然グループが証明した手法を商品化", ACCENT),
    ]
    for i, (phase, desc, color) in enumerate(phases):
        lx = Inches(0.3 + i * 4.3)
        add_rect(slide, lx, Inches(1.8), Inches(4.0), Inches(1.6), color)
        multi_txt(slide, phase.split("\n"), lx, Inches(1.85), Inches(4.0), Inches(0.7),
                  size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        multi_txt(slide, desc.split("\n"), lx, Inches(2.6), Inches(4.0), Inches(0.7),
                  size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # プラットフォーム別施策
    platforms = [
        ("YouTube", MPL_RED,
         ["然グループ公式チャンネル",
          "仕込み・料理動画・代表ブログ",
          "発酵・出汁の解説コンテンツ",
          "→ チャンネル登録→来店転換"]),
        ("Instagram", MPL_GOLD,
         ["ブランド別アカウント＋統合アカウント",
          "出汁の湯気・炭火の煙=映えコンテンツ",
          "ストーリーズで日常・ビハインドシーン",
          "→ フォロワーをリザーブに誘導"]),
        ("TikTok", MPL_NAVY,
         ["短尺料理動画・発酵解説",
          "インバウンド向け多言語対応",
          "ハッシュタグ戦略でバズ狙い",
          "→ 若年層・外国人の来店促進"]),
        ("Google / OTA", MPL_GREEN,
         ["Googleマップ口コミ管理",
          "食べログ・ぐるなびSEO",
          "Booking.com / Airbnb体験登録",
          "→ インバウンド流入確保"]),
    ]
    for i, (name, color, items) in enumerate(platforms):
        lx = Inches(0.3 + i * 3.25)
        add_rect(slide, lx, Inches(3.65), Inches(3.0), Inches(0.42), RGBColor(*bytes.fromhex(color[1:])))
        txt(slide, name, lx, Inches(3.68), Inches(3.0), Inches(0.35),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        multi_txt(slide, items, lx + Inches(0.1), Inches(4.12), Inches(2.8), Inches(1.8),
                  size=10, color=TEXT_COLOR)

    txt(slide, "開業前コンサル費用 300万円を資金調達に組み込み済み → Y1からSNS集客を最大化",
        Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.55),
        size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_group_annual_pl(prs):
    slide = blank(prs)
    add_header_band(slide, "グループ 事業計画 P&L  ─  5カ年損益計画")

    # 年次PL表
    hdrs = ["", "Y1 2026", "Y2 2027", "Y3 2028", "Y4 2029", "Y5 2030"]
    rows = [
        ["売上高（万円）",        "14,000",  "36,000",  "72,000",  "120,000", "200,000"],
        ["　売上原価（31%）",     "▲4,340",  "▲11,160", "▲22,320", "▲37,200", "▲62,000"],
        ["売上総利益",            "9,660",   "24,840",  "49,680",  "82,800",  "138,000"],
        ["　粗利率",              "69%",     "69%",     "69%",     "69%",     "69%"],
        ["　人件費",              "▲3,500",  "▲8,500",  "▲16,000", "▲26,000", "▲46,000"],
        ["　家賃合計",            "▲1,500",  "▲3,600",  "▲7,200",  "▲11,000", "▲17,000"],
        ["　その他（光熱費等）",  "▲2,400",  "▲5,600",  "▲11,600", "▲18,000", "▲35,000"],
        ["EBITDA",                "2,260",   "7,140",   "14,880",  "27,800",  "40,000"],
        ["EBITDA率",              "16%",     "20%",     "21%",     "23%",     "20%"],
    ]

    def highlight_rows(tbl, row_indices, color):
        for ri in row_indices:
            for ci in range(len(hdrs)):
                tbl.cell(ri+1, ci).fill.solid()
                tbl.cell(ri+1, ci).fill.fore_color.rgb = color

    n_rows = len(rows) + 1
    n_cols = len(hdrs)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.1), Inches(12.7), Inches(4.5))
    tbl = tbl_shape.table

    col_widths = [Inches(2.5)] + [Inches(2.04)] * 5
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(hdrs):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = FONT_NAME
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

    highlight_color_map = {
        2: RGBColor(0xE8, 0xF4, 0xEE),
        7: RGBColor(0xC4, 0x9A, 0x22),
    }
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri == 7:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            elif ri == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xEE)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_NAME
            run.font.size = Pt(10)
            run.font.bold = (ri in [2, 7])
            run.font.color.rgb = WHITE if ri == 7 else TEXT_COLOR

    # グラフ
    fig = fig_annual_pl_chart()
    embed_figure(slide, fig, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.4))


def slide_cashflow(prs):
    slide = blank(prs)
    add_header_band(slide, "キャッシュフロー計画")

    fig = fig_cashflow_chart()
    embed_figure(slide, fig, Inches(0.3), Inches(1.1), Inches(7.5), Inches(3.5))

    hdrs_cf = ["", "調達時", "Y1末", "Y2末", "Y3末", "Y4末", "Y5末"]
    rows_cf = [
        ["営業CF（万円）",    "─",      "+2,600", "+7,500",  "+15,000", "+28,000", "+48,000"],
        ["投資CF（万円）",    "▲3,800", "▲1,000", "▲5,000",  "▲8,000",  "▲10,000", "▲5,000"],
        ["純CF（万円）",      "+3,800", "+1,600", "+2,500",  "+7,000",  "+18,000", "+43,000"],
        ["累計CF（万円）",    "3,800",  "5,400",  "7,900",   "14,900",  "32,900",  "75,900"],
    ]
    table(slide, hdrs_cf, rows_cf, Inches(7.8), Inches(1.1), Inches(5.2), Inches(2.5),
          hdr_color=ACCENT, font_size=10)

    points = [
        "資金調達3,800万を起点にY1末で累計CF黒字転換（返済原資確保）",
        "Y2: FC展開への投資5,000万は営業CF7,140万の範囲内でまかなえる",
        "Y3以降: 投資超過分もFCロイヤルティ等のストック収益で補填、累計CFが急拡大",
        "Y5末: 累計CF 7.5億円 → 配当・EXIT・追加投資の原資として十分な水準",
    ]
    section_box(slide, "CFのポイント", points, Inches(7.8), Inches(3.8), Inches(5.2), ACCENT)

    # 月次Y1イメージ
    add_rect(slide, Inches(0.3), Inches(4.8), Inches(7.3), Inches(2.4), LIGHT_BG)
    txt(slide, "■ Y1 月次CF イメージ（肉酒場然・稼働済み前提）",
        Inches(0.4), Inches(4.85), Inches(7.1), Inches(0.4), size=11, bold=True, color=ACCENT)
    multi_txt(slide,
        ["月1〜3：肉酒場然のみ稼働 → 月次CF +100〜150万（立ち上げ費用除く）",
         "月4〜6：だし然 開業 → 月次CF +200〜280万（稼働率70%）",
         "月7〜：蕎麦然 開業 → 月次CF +350〜430万（3ブランド揃い踏み）",
         "Y1末時点：累計CF（資金調達後）約5,400万 → BEP到達"],
        Inches(0.4), Inches(5.3), Inches(7.1), Inches(1.8),
        size=10.5, color=TEXT_COLOR)


def slide_08_competitive(prs):
    slide = blank(prs)
    add_header_band(slide, "競合優位性  ─  なぜ然グループが勝てるか")

    fig = fig_competitive_map()
    embed_figure(slide, fig, Inches(0.3), Inches(1.1), Inches(5.3), Inches(4.5))

    advantages = [
        ("① 出汁・発酵という普遍的コア",
         "トレンドに流されない「日本食の本質」を軸に\n10年・20年後も価値が変わらない"),
        ("② 価格帯ポートフォリオ",
         "不況時は蕎麦然・ランチ中心で守り\n好況・インバウンドはだし然夜・肉酒場で攻める"),
        ("③ FC×直営の二重エンジン",
         "直営で利益を確保しながらFCでスケール\n飲食成功パターンを確実に踏む"),
        ("④ SNS情報優位",
         "毎日のトレンドリサーチを内製→外販できる\n唯一の飲食グループ"),
        ("⑤ インバウンド確保",
         "だし然がインバウンド向け体験型ディナーを担当\n円安・訪日外国人増加の波を確実に取込む"),
    ]
    for i, (title, body) in enumerate(advantages):
        lx = Inches(5.9)
        ty = Inches(1.1) + i * Inches(1.25)
        add_rect(slide, lx, ty, Inches(7.1), Inches(1.15), LIGHT_BG)
        add_rect(slide, lx, ty, Inches(0.08), Inches(1.15), ACCENT)
        txt(slide, title, lx + Inches(0.18), ty + Inches(0.05),
            Inches(6.8), Inches(0.4), size=12, bold=True, color=ACCENT)
        txt(slide, body, lx + Inches(0.18), ty + Inches(0.48),
            Inches(6.8), Inches(0.6), size=10, color=TEXT_COLOR)

    txt(slide, "調査対象：みなとや・ひまり堂・ひまり商店・びゃく・三谷・なかよし（恵比寿/新橋/学芸大学）→ いずれも「こだわり×展開力」の両立なし",
        Inches(0.3), Inches(6.1), Inches(12.7), Inches(0.4),
        size=9, color=TEXT_COLOR, align=PP_ALIGN.CENTER)
    txt(slide, "→ 「素材へのこだわり × FC展開力」が高水準で共存する競合不在のゾーン",
        Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5),
        size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_competitor_analysis(prs):
    """競合分析スライド①：6店比較表"""
    slide = blank(prs)
    add_header_band(slide, "競合分析  ─  ベンチマーク6店舗の詳細比較")

    # 説明文
    txt(slide, "調査エリア：恵比寿・新橋・学芸大学（いずれも然グループの出店候補商圏）",
        Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
        size=11, color=TEXT_COLOR, italic=True)

    # 比較表
    hdrs = ["店舗名", "エリア", "業態・コンセプト", "客単価", "店舗数/展開",
            "こだわり度", "展開可能性", "然グループとの関係"]
    rows = [
        ["みなとや",         "大手町",    "老舗和定食・ビジネス層向け",        "¥1,200",  "単店",         "★★★☆☆", "★★☆☆☆", "だし然の客層が重なる"],
        ["ひまり堂",         "恵比寿",    "こだわり出汁・和食カフェ",          "¥1,800",  "単店",         "★★★★☆", "★☆☆☆☆", "だし然の直接競合"],
        ["ひまり商店",       "新橋",      "居酒屋・サラリーマン向け",          "¥3,500",  "単店",         "★★☆☆☆", "★★★☆☆", "肉酒場然と客層が重なる"],
        ["びゃく",           "学芸大学",  "高級和食・職人こだわり",            "¥8,000",  "単店",         "★★★★★", "★☆☆☆☆", "価格帯上位・直接競合薄"],
        ["三谷",             "学芸大学",  "本格蕎麦・職人系単店",              "¥2,500",  "単店",         "★★★★☆", "★☆☆☆☆", "蕎麦然の品質ベンチ"],
        ["なかよし",         "恵比寿",    "カジュアル和食・女性向け",          "¥2,200",  "単店",         "★★★☆☆", "★★☆☆☆", "だし然ランチ帯と競合"],
    ]
    table(slide, hdrs, rows,
          Inches(0.3), Inches(1.5), Inches(12.7), Inches(4.5),
          hdr_color=ACCENT, font_size=10)

    # 凡例
    txt(slide, "※ こだわり度・展開可能性は5段階評価（★多い＝高い）　客単価は推計値",
        Inches(0.4), Inches(6.1), Inches(12.0), Inches(0.35),
        size=9, color=TEXT_COLOR, italic=True)

    txt(slide, "→ 高こだわり×高展開力を両立している競合はゼロ  ─  然グループの独占ゾーン",
        Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5),
        size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def fig_competitor_radar():
    """競合レーダーチャート"""
    import matplotlib.patches as mpatches
    categories = ["こだわり度", "FC展開力", "価格競争力", "SNS発信力", "ブランド統一"]
    N = len(categories)
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]

    # 各競合のスコア（0〜5）
    data = {
        "だし 然":   [4.0, 3.5, 3.5, 4.5, 5.0],
        "肉酒場 然": [4.5, 4.2, 3.0, 4.0, 5.0],
        "蕎麦 然":   [3.5, 4.8, 4.5, 3.5, 5.0],
        "ひまり堂":  [4.2, 1.0, 2.5, 2.0, 1.5],
        "びゃく":    [5.0, 1.0, 1.0, 1.5, 1.0],
        "三谷":      [4.5, 1.0, 2.0, 1.0, 1.0],
    }
    colors_map = {
        "だし 然": MPL_GOLD,
        "肉酒場 然": MPL_RED,
        "蕎麦 然": MPL_NAVY,
        "ひまり堂": "#aaaaaa",
        "びゃく": "#888888",
        "三谷": "#666666",
    }

    fig, ax = plt.subplots(figsize=(5.2, 4.5), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor(MPL_LIGHT)
    ax.set_facecolor(MPL_LIGHT)

    for name, vals in data.items():
        vals_plot = vals + vals[:1]
        lw = 2.5 if name in ["だし 然", "肉酒場 然", "蕎麦 然"] else 1.2
        ls = "-" if name in ["だし 然", "肉酒場 然", "蕎麦 然"] else "--"
        alpha = 0.25 if name in ["だし 然", "肉酒場 然", "蕎麦 然"] else 0.0
        ax.plot(angles, vals_plot, linewidth=lw, linestyle=ls,
                color=colors_map[name], label=name)
        ax.fill(angles, vals_plot, color=colors_map[name], alpha=alpha)

    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontfamily="IPAGothic", fontsize=9)
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(["1", "2", "3", "4", "5"], fontsize=7, color="gray")
    ax.grid(color="gray", linestyle="--", linewidth=0.5, alpha=0.5)
    ax.set_title("競合5軸比較（レーダー）", fontfamily="IPAGothic",
                 fontsize=11, color=MPL_GREEN, fontweight="bold", pad=15)
    ax.legend(loc="upper right", bbox_to_anchor=(1.35, 1.1),
              prop={"family": "IPAGothic", "size": 8})
    fig.tight_layout()
    return fig


def slide_competitor_strategy(prs):
    """競合分析スライド②：差別化戦略サマリー"""
    slide = blank(prs)
    add_header_band(slide, "競合分析  ─  差別化戦略と勝ち筋")

    # レーダーチャート（左）
    fig = fig_competitor_radar()
    embed_figure(slide, fig, Inches(0.3), Inches(1.05), Inches(5.8), Inches(5.2))

    # 右側：差別化ポイントカード
    insights = [
        ("ひまり堂（恵比寿）との差",
         "• こだわり度は同等だが展開力が決定的に違う\n• 然グループはFC化で1店の成功を30店に複製できる\n• ひまり堂は単店の「職人店」─ スケールしない"),
        ("びゃく・三谷との差",
         "• 職人系の最高品質だが客単価¥8,000超で市場が狭い\n• 蕎麦然は¥2,500帯で「日常使いできるこだわり蕎麦」\n• 立地も8坪小型で回転率×FC数で勝負"),
        ("ひまり商店・なかよしとの差",
         "• 価格帯は近いが「素材・発酵」の軸が弱い\n• 然グループは出汁・麹・発酵を共通コアとして全ブランドが連動\n• SNS情報力で口コミ・インバウンドを自力で取込む"),
        ("然グループ固有の強み（競合ゼロ）",
         "• 3ブランド×FC展開：1業態の不振を他で補うリスクヘッジ\n• 素材へのこだわり + FC展開力の両立 = 競合不在ゾーン\n• SNS内製 = 月次マーケ費ゼロで情報優位を維持"),
    ]

    colors_list = [GOLD, NAVY, DARK_RED, ACCENT]
    for i, (title, body) in enumerate(insights):
        ty = Inches(1.1) + i * Inches(1.52)
        add_rect(slide, Inches(6.3), ty, Inches(6.7), Inches(1.42), LIGHT_BG)
        add_rect(slide, Inches(6.3), ty, Inches(0.1), Inches(1.42), colors_list[i])
        txt(slide, title, Inches(6.5), ty + Inches(0.06),
            Inches(6.4), Inches(0.38), size=11, bold=True, color=colors_list[i])
        txt(slide, body, Inches(6.5), ty + Inches(0.44),
            Inches(6.4), Inches(0.95), size=9.5, color=TEXT_COLOR)

    txt(slide, "結論：然グループは「こだわり × スケール」の両立という、現状競合が存在しないポジションを確立する",
        Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.5),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def fig_business_model_comparison():
    """ビジネスモデル比較：単店職人 vs 然グループ FC"""
    fig, axes = plt.subplots(1, 2, figsize=(9.5, 3.8))
    fig.patch.set_facecolor(MPL_LIGHT)

    # 左：競合（単店職人）モデルのキャッシュフローイメージ
    ax1 = axes[0]
    ax1.set_facecolor(MPL_LIGHT)
    years = ["1年目", "3年目", "5年目", "10年目"]
    rival_cf = [100, 150, 160, 170]   # 単店頭打ち
    ax1.bar(range(len(years)), rival_cf, color="#aaaaaa", alpha=0.8)
    ax1.set_xticks(range(len(years)))
    ax1.set_xticklabels(years, fontfamily="IPAGothic", fontsize=9)
    ax1.set_title("競合（単店職人モデル）", fontfamily="IPAGothic", fontsize=10,
                  color="#666666")
    ax1.set_ylabel("月次売上（万円）", fontfamily="IPAGothic", fontsize=9)
    ax1.set_ylim(0, 600)
    for i, v in enumerate(rival_cf):
        ax1.text(i, v + 8, f"{v}万", ha="center", fontfamily="IPAGothic", fontsize=9, color="#555")
    ax1.annotate("※ 単店で頭打ち\nスケールしない", xy=(3, 170), xytext=(2.0, 340),
                 arrowprops=dict(arrowstyle="->", color="#cc4444"),
                 fontfamily="IPAGothic", fontsize=9, color="#cc4444")
    ax1.spines[["top", "right"]].set_visible(False)

    # 右：然グループ FC展開モデル
    ax2 = axes[1]
    ax2.set_facecolor(MPL_LIGHT)
    zen_cf = [120, 300, 600, 1700]  # FC店舗数に応じてスケール
    ax2.bar(range(len(years)), zen_cf, color=[MPL_GOLD, MPL_GOLD, MPL_GREEN, MPL_GREEN], alpha=0.85)
    ax2.set_xticks(range(len(years)))
    ax2.set_xticklabels(years, fontfamily="IPAGothic", fontsize=9)
    ax2.set_title("然グループ（FC展開モデル）", fontfamily="IPAGothic", fontsize=10,
                  color=MPL_GREEN)
    ax2.set_ylabel("月次グループ売上（万円）", fontfamily="IPAGothic", fontsize=9)
    ax2.set_ylim(0, 2200)
    for i, v in enumerate(zen_cf):
        ax2.text(i, v + 30, f"{v}万", ha="center", fontfamily="IPAGothic", fontsize=9,
                 color=MPL_GREEN)
    ax2.annotate("FC×複数ブランドで\n10倍以上にスケール", xy=(3, 1700), xytext=(1.8, 1900),
                 arrowprops=dict(arrowstyle="->", color=MPL_GREEN),
                 fontfamily="IPAGothic", fontsize=9, color=MPL_GREEN)
    ax2.spines[["top", "right"]].set_visible(False)

    fig.tight_layout(pad=2.0)
    return fig


def slide_differentiation_deep(prs):
    """差別化深掘りスライド：FCモデル × 発酵コアの独自性"""
    slide = blank(prs)
    add_header_band(slide, "差別化の本質  ─  FCモデル × 発酵・出汁コアが生む参入障壁")

    # ── 上段：ビジネスモデル比較チャート ──
    fig = fig_business_model_comparison()
    embed_figure(slide, fig, Inches(0.3), Inches(1.05), Inches(9.8), Inches(3.4))

    # 右上：FC優位性の数値サマリー
    fc_points = [
        ("1店→55店への複製力", "成功レシピをFC化で横展開\n競合は単店の技術・関係性に依存"),
        ("ロイヤルティ収入", "売上の3〜4%が直営費用ゼロで入る\n店舗数×ロイヤルティ＝純粋な利益エンジン"),
        ("リスク分散", "3ブランド×FC＝景気・トレンドに左右されない\n単店競合は1業態の不振で即ダメージ"),
    ]
    for i, (title, body) in enumerate(fc_points):
        ty = Inches(1.1) + i * Inches(1.1)
        add_rect(slide, Inches(10.3), ty, Inches(2.9), Inches(1.0), LIGHT_BG)
        add_rect(slide, Inches(10.3), ty, Inches(0.08), Inches(1.0), GOLD)
        txt(slide, title, Inches(10.45), ty + Inches(0.04),
            Inches(2.7), Inches(0.35), size=9.5, bold=True, color=GOLD)
        txt(slide, body, Inches(10.45), ty + Inches(0.38),
            Inches(2.7), Inches(0.58), size=8.5, color=TEXT_COLOR)

    # ── 下段：発酵・出汁コア ──
    add_rect(slide, Inches(0.3), Inches(4.6), Inches(12.7), Inches(0.38), ACCENT)
    txt(slide, "素材・発酵コアの独自性  ─  競合が真似できない「知識資産」",
        Inches(0.5), Inches(4.63), Inches(12.3), Inches(0.35),
        size=12, bold=True, color=WHITE)

    ferment_items = [
        ("出汁の科学的再現性",
         "昆布・鰹・椎茸の配合レシピを数値化\n職人の勘ではなく「再現できる出汁」\n→ FC加盟店でも同品質を担保できる",
         GOLD),
        ("麹・発酵の内製ノウハウ",
         "麹仕入れKg/¥1,300（税抜）から\n塩麹・醤油麹・発酵調味料を内製化\n→ 原価コントロールと差別化を同時実現",
         NAVY),
        ("競合に真似されない理由",
         "ひまり堂・三谷・びゃくは職人個人の感覚に依存\n然グループは「再現可能なレシピ資産」として蓄積\n→ 人が替わっても品質が変わらない = FC化の前提",
         ACCENT),
        ("グループ横断のシナジー",
         "だし然の出汁知識 → 蕎麦然のつゆに応用\n麹・発酵は全ブランドの共通コア食材\n→ R&Dコストをグループで共有、競合はブランドごとに個別対応",
         DARK_RED),
    ]

    for i, (title, body, color) in enumerate(ferment_items):
        lx = Inches(0.3) + i * Inches(3.2)
        add_rect(slide, lx, Inches(5.05), Inches(3.05), Inches(2.1), LIGHT_BG)
        add_rect(slide, lx, Inches(5.05), Inches(3.05), Inches(0.07), color)
        txt(slide, title, lx + Inches(0.1), Inches(5.13),
            Inches(2.85), Inches(0.35), size=10, bold=True, color=color)
        txt(slide, body, lx + Inches(0.1), Inches(5.5),
            Inches(2.85), Inches(1.58), size=8.5, color=TEXT_COLOR)

    txt(slide, "→ 「再現可能なこだわり」という矛盾を解決しているのが然グループだけ",
        Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.35),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_09_5year(prs):
    slide = blank(prs)
    add_header_band(slide, "5カ年数値計画")

    # 成長グラフ
    fig = fig_growth_chart()
    embed_figure(slide, fig, Inches(0.3), Inches(1.1), Inches(7.5), Inches(4.2))

    # 数値表（右側）
    hdrs = ["年度", "店舗数", "月商", "年商"]
    rows = [
        ["Y1 2026", "3店", "1,200万", "1.4億"],
        ["Y2 2027", "9店", "3,000万", "3.6億"],
        ["Y3 2028", "22店", "6,000万", "7.2億"],
        ["Y4 2029", "35店", "10,000万", "12億"],
        ["Y5 2030", "55店", "17,000万", "20億"],
    ]
    table(slide, hdrs, rows, Inches(8.0), Inches(1.1), Inches(5.0), Inches(3.3), font_size=12)

    add_rect(slide, Inches(8.0), Inches(4.6), Inches(5.0), Inches(0.7), GOLD)
    txt(slide, "Year 5 目標", Inches(8.1), Inches(4.65), Inches(4.8), Inches(0.3),
        size=10, color=WHITE, bold=True)
    txt(slide, "年商 20億円  /  55店舗  /  EBITDA 4億円",
        Inches(8.1), Inches(4.95), Inches(4.8), Inches(0.35),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # マイルストーン
    milestones = [
        ("Y1", "だし然・蕎麦然1号店開業", GOLD),
        ("Y2", "FC展開開始・肉酒場然2号店", DARK_RED),
        ("Y3", "関西進出・SNS事業部外販", NAVY),
        ("Y5", "55店舗・年商20億・EXIT準備", ACCENT),
    ]
    for i, (year, text, color) in enumerate(milestones):
        lx = Inches(0.3 + i * 3.2)
        add_rect(slide, lx, Inches(5.55), Inches(3.0), Inches(1.6), color)
        txt(slide, year, lx, Inches(5.6), Inches(3.0), Inches(0.45),
            size=18, bold=True, color=GOLD if color != GOLD else WHITE,
            align=PP_ALIGN.CENTER)
        txt(slide, text, lx + Inches(0.1), Inches(6.1), Inches(2.8), Inches(0.9),
            size=10, color=WHITE, align=PP_ALIGN.CENTER)


def slide_10_revenue(prs):
    slide = blank(prs)
    add_header_band(slide, "収益モデルサマリー  ─  グループ損益と成長ストック")

    # 左上：グループPL比較グラフ
    fig = fig_group_pl()
    embed_figure(slide, fig, Inches(0.3), Inches(1.1), Inches(6.2), Inches(3.1))

    # 右：P&L表
    hdrs = ["ブランド", "月次売上/店", "利益率", "1店月次利益", "目標店舗", "グループ月次利益"]
    rows = [
        ["だし 然", "960万円", "15%", "144万円", "20店", "2,880万円"],
        ["肉酒場 然", "1,200万円", "17%", "200万円", "5店", "1,000万円"],
        ["蕎麦 然（小型）", "350万円", "20%", "70万円", "30店", "2,100万円"],
        ["SNS事業部", "─", "─", "83万円", "─", "83万円"],
        ["合計（安定期）", "─", "─", "─", "55店", "6,063万円〜"],
    ]
    table(slide, hdrs, rows, Inches(6.5), Inches(1.1), Inches(6.5), Inches(3.1), font_size=10)

    # ポイント
    points = [
        "Y1 グループ月次売上：約2,510万円（利益率約17%、月次利益約418万円）",
        "FCロイヤルティはFC40店規模で月1,296万円のストック収益（安定成長）",
        "成長フェーズ（Y3〜）で利益率が改善：規模の経済＋仕入共通化",
        "直営利益＋FCロイヤルティの二重収益でリスク分散",
    ]
    section_box(slide, "収益構造のポイント", points,
                Inches(0.3), Inches(4.4), Inches(12.7), ACCENT)


def slide_11_funding(prs):
    slide = blank(prs)
    add_header_band(slide, "資金調達計画")

    add_rect(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(1.3), ACCENT)
    txt(slide, "第1回 調達目標", Inches(0.4), Inches(1.15), Inches(4.3), Inches(0.4),
        size=13, color=WHITE)
    txt(slide, "3,800万円", Inches(0.4), Inches(1.55), Inches(4.3), Inches(0.75),
        size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.1), Inches(1.1), Inches(7.9), Inches(1.3), LIGHT_BG)
    txt(slide, "調達方法：日本政策金融公庫 ＋ エンジェル投資家",
        Inches(5.2), Inches(1.35), Inches(7.7), Inches(0.5),
        size=15, bold=True, color=ACCENT)
    txt(slide, "第2回（Y2〜3）：1〜2億円  /  VC・事業会社提携  /  FC加速・関西展開",
        Inches(5.2), Inches(1.85), Inches(7.7), Inches(0.4),
        size=11, color=TEXT_COLOR)

    hdrs = ["用途", "金額", "内訳"]
    rows = [
        ["だし 然 1号店 開業費", "2,000万円", "内装900万＋厨房450万＋敷金（60万×3ヶ月）＋運転資金"],
        ["蕎麦 然 1号店 開業費（8坪）", "800万円", "内装250万＋設備200万＋敷金（45万×6ヶ月）＋運転資金"],
        ["SNSコンサル・スタートダッシュ", "300万円", "外部コンサル起用・開業前3〜6ヶ月集中投資"],
        ["人材採用・研修費", "300万円", "店長・料理長採用・マニュアル・研修体制構築"],
        ["運転資金・予備費", "400万円", "グループ管理・広告・オープン販促・予備"],
        ["合計", "3,800万円", ""],
    ]
    table(slide, hdrs, rows, Inches(0.3), Inches(2.6), Inches(12.7), Inches(3.2),
          hdr_color=ACCENT, font_size=11)

    multi_txt(slide,
        ["■ 投資家へのリターンイメージ",
         "・Y1〜2：初期投資回収フェーズ。各店舗の安定稼働と収益基盤の構築に集中",
         "・Y3以降：FCスケールで利益率が改善。EBITDA1.8億円（月商6,000万）で成長が加速",
         "・Y5：年商20億円・EBITDA4億円  →  EXIT（M&A or IPO）を視野に",
         "・FCロイヤルティ収入はストック型のため、EXIT評価倍率向上に寄与"],
        Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.6),
        size=11, color=TEXT_COLOR)


def slide_12_risk(prs):
    slide = blank(prs)
    add_header_band(slide, "リスクと対策")

    hdrs = ["リスク", "影響度", "発生確率", "対策"]
    rows = [
        ["物件取得競争", "高", "中",
         "不動産仲介との早期リレーション・複数候補常時確保"],
        ["FC加盟店の品質低下", "高", "低",
         "SV制度・月次研修・本部仕入れで品質統一・契約解除条項"],
        ["インバウンド需要の変動", "中", "中",
         "ランチ（国内需要）でベース売上確保・円安依存排除"],
        ["人材採用難", "中", "高",
         "SNSブランディングで「働きたい会社」化・利益還元・待遇改善"],
        ["出汁素材の調達リスク", "低", "低",
         "複数サプライヤーと年間契約・輸入品との組合せで安定化"],
        ["競合参入", "中", "中",
         "先行出店・ブランド確立・FC網を早期に張る先行者利益戦略"],
    ]
    table(slide, hdrs, rows, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5),
          hdr_color=ACCENT, font_size=11)


def slide_13_contact(prs):
    slide = blank(prs, dark=True)
    txt(slide, "然グループ", Inches(1), Inches(1.2), Inches(11), Inches(1.5),
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, "ZEN GROUP", Inches(1), Inches(2.7), Inches(11), Inches(0.5),
        size=16, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    add_rect(slide, Inches(3), Inches(3.5), Inches(7), Inches(0.04), GOLD)
    txt(slide, "食の全シーンに、然を。",
        Inches(1), Inches(3.7), Inches(11), Inches(0.6),
        size=20, color=GOLD, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3), Inches(4.5), Inches(7), Inches(2.0),
             RGBColor(0x1E, 0x35, 0x2A))
    txt(slide, "然グループ 事業企画室",
        Inches(3.1), Inches(4.6), Inches(6.8), Inches(0.5),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, "Email：＿＿＿＿＿＿＿＿＿＿＿＿＿",
        Inches(3.1), Inches(5.15), Inches(6.8), Inches(0.45),
        size=14, color=LIGHT_BG, align=PP_ALIGN.CENTER)
    txt(slide, "Tel：＿＿＿＿＿＿＿＿＿",
        Inches(3.1), Inches(5.6), Inches(6.8), Inches(0.45),
        size=14, color=LIGHT_BG, align=PP_ALIGN.CENTER)
    txt(slide, "本資料は機密情報を含みます。無断転載・配布を禁じます。",
        Inches(1), Inches(6.8), Inches(11), Inches(0.4),
        size=9, color=GOLD, align=PP_ALIGN.CENTER)


# ─── メイン ─────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_vision(prs)
    slide_03_market(prs)
    slide_04_overview(prs)
    slide_05_dashi(prs)
    slide_06_niku(prs)
    slide_07_soba(prs)
    slide_08_competitive(prs)
    slide_competitor_analysis(prs)
    slide_competitor_strategy(prs)
    slide_differentiation_deep(prs)
    slide_sns_marketing(prs)
    slide_09_5year(prs)
    slide_10_revenue(prs)
    slide_group_annual_pl(prs)
    slide_cashflow(prs)
    slide_11_funding(prs)
    slide_12_risk(prs)
    slide_13_contact(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"Saved: {OUTPUT_FILE}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
